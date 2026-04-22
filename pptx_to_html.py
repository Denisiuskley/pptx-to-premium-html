import os
import json
import re
import io
import html
import math
import logging
import hashlib
import time
import base64
import mimetypes
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import requests
import xml.etree.ElementTree as ET

try:
    from lxml import etree

    HAS_LXML = True
except ImportError:
    HAS_LXML = False
    etree = None

# Office Math Markup Language (OMML) namespace
OMML_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/math}"
MATHML_NS = "{http://www.w3.org/1998/Math/MathML}"

# Константы порогов для динамической типографии и layout
TEXT_LEN_HEAVY = 500  # Если текст длиннее, делаем левую колонку шире
TEXT_LEN_LIGHT = 200  # Если текст короче, делаем правую колонку (изображение) шире
TEXT_LEN_CONDENSED = 800  # Порог для перехода на сжатый шрифт (для узких колонок)
TEXT_LEN_TIGHT = 1200  # Порог для перехода на очень сжатый шрифт (для узких колонок)
ASPECT_WIDE = 1.8  # Соотношение сторон, при котором изображение считаются широким
ASPECT_TALL = 0.7  # Соотношение сторон, при котором изображение считается высоким
CAPTION_H_DIST_FACTOR = 0.8  # Максимальное горизонтальное расстояние между подписью и изображением (относительно ширины изображения)
CAPTION_V_GAP_MIN = (
    -5000
)  # Минимальный вертикальный зазор между подписью и изображением
CAPTION_V_GAP_MAX = 30000  # Максимальный вертикальный зазор

logging.basicConfig(level=logging.INFO, format="[%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

# Базовые пути
BASE_DIR = Path(__file__).parent.resolve()


# ==============================================================================
# ВСТРОЕННЫЙ HTML-ШАБЛОН (все стили и разметка вшиты в скрипт)
# ==============================================================================
BASE_HTML_TEMPLATE = """<!DOCTYPE html>
<html class="no-js" lang="ru">

<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Доклад: {speaker_name} (ПНИПУ)</title>
    <link rel="stylesheet" href="libs/fonts/fonts.css">
    <script src="libs/gsap/gsap.min.js" defer></script>
    <script src="libs/lucide/lucide.min.js" defer></script>
    <!-- Local MathJax configuration -->
    <script>
        window.MathJax = {
            options: {
                enableExplorer: false, // Отключаем Explorer, который может перекрывать формулы желтыми плашками
                enableAssistiveMml: true
            },
            loader: {
                load: ['[tex]/ams']
            },
            svg: {
                fontCache: 'global',
                scale: 1.2
            },
            startup: {
                ready: () => {
                    console.log('MathJax is ready (SVG mode)!');
                    MathJax.startup.defaultReady();
                    MathJax.startup.promise.then(() => {
                        // Анимация появления формул после рендеринга
                        gsap.to('.formula-container, .formula-block', {
                            opacity: 1,
                            y: 0,
                            duration: 0.8,
                            stagger: 0.1,
                            ease: 'power2.out'
                        });
                    });
                }
            }
        };
    </script>
    <script src="libs/mathjax/tex-mml-svg.js" defer></script>
        <style>
            :root {
                --bg-deep: #05080e;
                --bg-card: rgba(15, 23, 42, 0.7);
                --accent: #00f2ff;
                --accent-soft: rgba(0, 242, 255, 0.1);
                --text-main: #f8fafc;
                --text-dim: #cbd5e1;
                --glass-border: rgba(255, 255, 255, 0.1);
                --gradient-accent: linear-gradient(135deg, #00f2ff 0%, #0066ff 100%);
                --glow-cyan: 0 0 20px rgba(0, 242, 255, 0.4);

                --logo-height: 160px;
                --logo-top: -1.0rem;
                --logo-right: 4rem;
                --header-height: 75px;
                --slide-padding-v: 2rem;
                --slide-padding-h: 4rem;
                --panel-padding: 2.0rem;

                --fs-main-title: 4.2rem;
                --fs-slide-title: 1.4rem;
                --fs-slide-num: 1.1rem;
                --fs-sub-heading: 1.8rem;
                --fs-text-main: 1.3rem;
                --fs-tag: 0.9rem;
                --fs-viz-caption: 1.1rem;
                --fs-viz-desc: 1.0rem;
                --fs-presenter-name: 1.5rem;
                --fs-presenter-info: 1.1rem;
                --fs-presenter-label: 0.8rem;
                --fs-research-year: 5rem;

                /* Таблицы */
                --fs-table-th: 1.05rem;
                --fs-table-td: 1.1rem;
                --fs-formula: 1.25rem;
                --gap-main: 2rem;
                --gap-items: 1.2rem;

                /* Adaptive Squeeze Factors */
                --squeeze-factor: 1;
                --base-gap: 0.5rem;
                --base-margin: 0.5rem;
                --base-lh: 1.6;

                /* Config-driven variables */
                --icon-size: 1.6rem;
                --bullet-icon-size: 1.6rem;
                --bullet-indent: 2rem;
                --bullet-border: 2px solid rgba(0, 242, 255, 0.15);
                --bullet-bg: rgba(0, 242, 255, 0.03);
                --formula-fallback-font: 'Roboto Mono', monospace;
                --formula-fallback-size: 1.1em;
            }

        * {
            margin: 0;
            padding: 0;
            box-sizing: border-box;
            -webkit-font-smoothing: antialiased;
        }

        body {
            background-color: var(--bg-deep);
            color: var(--text-main);
            font-family: 'Inter', sans-serif;
            overflow: hidden;
            height: 100vh;
        }

        .presentation {
            position: relative;
            width: 100vw;
            height: 100vh;
            display: flex;
            transition: transform 0.8s cubic-bezier(0.85, 0, 0.15, 1);
        }

        .slide {
            flex: 0 0 100vw;
            height: 100vh;
            display: flex;
            flex-direction: column;
            padding: var(--slide-padding-v) var(--slide-padding-h);
            position: relative;
            overflow: hidden;
        }

        .bg-element {
            position: absolute;
            z-index: -1;
            filter: blur(80px);
            opacity: 0.15;
            border-radius: 50%;
        }

        .bg-1 {
            width: 600px;
            height: 600px;
            background: var(--gradient-accent);
            top: -200px;
            right: -100px;
        }

        .bg-2 {
            width: 400px;
            height: 400px;
            background: #ff00ea;
            bottom: -100px;
            left: -100px;
            opacity: 0.05;
        }

        .slide-header {
            display: flex;
            justify-content: flex-start;
            align-items: center;
            gap: 1.5rem;
            margin-bottom: var(--margin-header);
            border-bottom: 1px solid var(--glass-border);
            padding-bottom: 0.5rem;
            position: relative;
            height: var(--header-height);
        }

        .logo-container {
            position: absolute;
            top: var(--logo-top);
            right: var(--logo-right);
            height: var(--logo-height);
            display: flex;
            align-items: center;
            z-index: 10;
        }

        .header-logo {
            height: 100%;
            width: auto;
            filter: drop-shadow(0 0 20px rgba(0, 242, 255, 0.5));
            opacity: 0.9;
        }

        .hide-title .slide-title {
            display: none;
        }

        .slide-title {
            font-family: 'Outfit', sans-serif;
            font-size: var(--fs-slide-title);
            letter-spacing: 0.05em;
            text-transform: uppercase;
            color: var(--accent);
            font-weight: 700;
        }

        .slide-number {
            font-family: 'Roboto Mono', monospace;
            font-size: var(--fs-slide-num);
            color: var(--text-dim);
            padding-right: 1rem;
            border-right: 1px solid var(--glass-border);
        }

        .slide-content-title {
            flex: 1;
            display: flex;
            flex-direction: column;
            justify-content: center;
            max-width: 1400px;
            margin: 0 auto;
            width: 100%;
        }

        .main-heading {
            font-family: 'Inter', sans-serif;
            font-size: var(--fs-main-title);
            font-weight: 800;
            line-height: 1.1;
            margin-bottom: 3rem;
            text-align: center;
            background: linear-gradient(to bottom right, #fff 50%, #94a3b8);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
        }

        .presenter-card {
            background: var(--bg-card);
            backdrop-filter: blur(12px);
            border: 1px solid var(--glass-border);
            padding: 2rem;
            border-right: 4px solid var(--accent);
            align-self: flex-start;
            width: 500px;
            position: relative;
            z-index: 20;
        }

        .presenter-label {
            color: var(--accent);
            font-size: var(--fs-presenter-label);
            text-transform: uppercase;
            letter-spacing: 0.2em;
            margin-bottom: 0.5rem;
            display: block;
        }

        .presenter-name {
            font-size: var(--fs-presenter-name);
            font-weight: 800;
            margin-bottom: 0.25rem;
            color: var(--text-main);
        }

        .presenter-info {
            color: var(--text-dim);
            font-size: var(--fs-presenter-info);
            line-height: 1.6;
            word-wrap: break-word;
            overflow-wrap: anywhere;
            max-width: 90%;
        }

        /* Presenter card positioned at left-bottom on intro slide */
        .slide.hide-title .presenter-card {
            position: absolute;
            bottom: var(--slide-padding-v);
            left: var(--slide-padding-h);
        }

        .slide-split {
            display: grid;
            grid-template-columns: 1.2fr 1.8fr;
            gap: var(--gap-main);
            flex: 1;
            min-height: 0;
        }

        .img-stack {
            display: flex;
            flex-direction: row;
            gap: var(--gap-main);
            align-items: stretch;
            height: 100%;
            min-height: 0;
            width: 100%;
        }

        .formula-block {
            margin: 1.5rem 0;
            text-align: center;
            background: rgba(0, 242, 255, 0.02);
            padding: 1rem;
            border-radius: 8px;
            border: 1px solid var(--glass-border);
        }

        /* --- Стили для адаптивной верстки (одна колонка рисунков - одиночные или в ряд) --- */
        .slide-split.layout-auto-width {
            grid-template-columns: minmax(25%, 1fr) auto !important;
        }
        .layout-auto-width .img-stack {
            width: auto;
            justify-content: flex-end;
        }
        .layout-auto-width .viz-item {
            width: auto;
            flex: none;
        }
        .layout-auto-width .viz-box {
            width: auto;
            max-width: 65vw;
        }
        .layout-auto-width .viz-box img {
            width: auto;
            height: 100%;
            object-fit: contain;
        }

        .viz-item {
            display: flex;
            flex-direction: column;
            gap: 1rem;
            flex: 1;
            min-height: 0;
        }

        .viz-caption {
            border-left: 4px solid var(--accent);
            padding: calc(0.5rem * var(--squeeze-factor)) calc(1rem * var(--squeeze-factor));
            background: linear-gradient(to right, var(--accent-soft), transparent);
            color: var(--text-main);
            font-size: var(--fs-viz-caption);
            font-weight: 600;
            line-height: 1.3;
        }

        .data-table {
            width: 100%;
            border-collapse: collapse;
            margin-top: 1rem;
            font-size: var(--fs-table-td);
            color: var(--text-main);
            font-family: 'Inter', sans-serif;
        }
        .data-table th, .data-table td {
            border: 1px solid var(--glass-border);
            padding: 0.25rem 0.5rem;
            text-align: left;
            vertical-align: top;
            line-height: 1.3;
            word-break: break-word;
            overflow-wrap: break-word;
        }
        .data-table th {
            background: var(--accent-soft);
            color: var(--accent);
            font-weight: 700;
            letter-spacing: 0.05em;
            font-size: var(--fs-table-th);
        }
        .data-table tr:nth-child(even) {
            background: rgba(255, 255, 255, 0.02);
        }
        .data-table tr:hover {
            background: rgba(0, 242, 255, 0.05);
        }

        .analytical-panel {
            background: var(--bg-card);
            backdrop-filter: blur(8px);
            border: 1px solid var(--glass-border);
            padding: var(--panel-padding);
            border-radius: 8px;
            display: flex;
            flex-direction: column;
            gap: calc(var(--base-gap) * var(--squeeze-factor));
            overflow-y: hidden;
            box-shadow: inset 0 1px 0 rgba(255,255,255,0.1), inset 1px 0 0 rgba(255,255,255,0.05), 0 8px 32px rgba(0,0,0,0.3);
            position: relative;
        }



        .section-tag {
            font-family: 'Roboto Mono', monospace;
            background: var(--accent-soft);
            color: var(--accent);
            padding: 4px 12px;
            font-size: var(--fs-tag);
            width: fit-content;
            border-radius: 4px;
        }

        .sub-heading {
            font-size: var(--fs-sub-heading);
            font-weight: 600;
            margin-bottom: calc(var(--base-margin) * var(--squeeze-factor));
        }

        .list-item {
            display: flex;
            gap: 1rem;
            margin-bottom: calc(var(--base-margin) * var(--squeeze-factor));
            align-items: flex-start;
        }

        .list-item i {
            margin-top: 4px;
            color: var(--accent);
            flex-shrink: 0;
            width: var(--icon-size);
            height: var(--icon-size);
        }

        /* Bullet list styling - distinct visual level */
        .list-item-bullet {
            display: flex;
            gap: 1rem;
            margin-bottom: calc(var(--base-margin) * var(--squeeze-factor));
            align-items: flex-start;
            padding-left: var(--bullet-indent);
            border-left: var(--bullet-border);
            background: var(--bullet-bg);
            border-radius: 0 8px 8px 0;
            padding: 0.5rem 1rem 0.5rem 0.5rem;
        }
        .list-item-bullet i {
            margin-top: 4px;
            color: var(--accent);
            flex-shrink: 0;
            width: var(--bullet-icon-size);
            height: var(--bullet-icon-size);
            opacity: 0.85;
        }

        .viz-container:has(.viz-item) {
            max-width: fit-content;
        }
         .viz-container:has(.viz-item) .viz-item {
             width: 100%;
         }

         .viz-item.wide {
             grid-column: span 2;
         }
         .viz-item.tall {
             grid-row: span 2;
         }

        .list-text {
            color: var(--text-dim);
            font-size: var(--fs-text-main);
            line-height: calc(var(--base-lh) * var(--squeeze-factor));
        }

        .list-text strong {
            color: var(--text-main);
        }

        /* Formula fallback styling */
        .formula-fallback {
            font-family: var(--formula-fallback-font, 'Roboto Mono'), monospace;
            font-size: var(--formula-fallback-size, 1.1em);
            white-space: nowrap;
            color: var(--accent);
            opacity: 0.9;
            display: inline-block;
            padding: 0.1rem 0.3rem;
            background: rgba(0, 242, 255, 0.05);
            border-radius: 4px;
        }

        .formula-container {
            display: inline-block;
            vertical-align: middle;
            margin: 0 0.3rem;
            opacity: 0;
            transform: translateY(5px);
        }

        .formula-block {
            background: var(--bg-card);
            backdrop-filter: blur(12px);
            border: 1px solid var(--glass-border);
            padding: 1rem 1.5rem;
            margin: 0.75rem 0;
            display: flex;
            justify-content: center;
            align-items: center;
            box-shadow: 0 12px 40px rgba(0,0,0,0.4);
            opacity: 0;
            transform: translateY(10px);
            transition: border-color 0.4s, box-shadow 0.4s;
        }

        .formula-block:hover {
            border-color: var(--accent);
            box-shadow: 0 0 30px rgba(0, 242, 255, 0.15);
        }

        /* MathJax SVG tuning */
        mjx-container[jax="SVG"] {
            color: var(--text-main);
            margin: 0 !important;
            padding: 0 !important;
        }

        .formula-container mjx-container[jax="SVG"] {
            display: inline-block !important;
            vertical-align: middle !important;
        }

        .viz-box {
            background: radial-gradient(circle at center, rgba(0, 242, 255, 0), transparent);
            border: 1px solid rgba(255, 255, 255, 0.1);
            border-radius: 12px;
            display: flex;
            align-items: center;
            justify-content: center;
            position: relative;
            overflow: hidden;
            height: 100%;
            width: 100%;
            min-height: 0;
            cursor: zoom-in;
            transition: all 0.5s cubic-bezier(0.4, 0, 0.2, 1);
        }

        .viz-box:hover {
            border-color: var(--accent);
            box-shadow:
                0 0 20px rgba(0, 242, 255, 0.2),
                inset 0 0 30px rgba(0, 242, 255, 0.15);
        }

        .viz-box::after {
            content: '';
            position: absolute;
            top: 0;
            left: 0;
            right: 0;
            bottom: 0;
            background-image:
                linear-gradient(rgba(0, 242, 255, 0.05) 1px, transparent 1px),
                linear-gradient(90deg, rgba(0, 242, 255, 0.05) 1px, transparent 1px);
            background-size: 20px 20px;
            opacity: 0;
            transition: opacity 0.4s;
            pointer-events: none;
        }

        .viz-box:hover::after {
            opacity: 1;
        }

        .viz-box img {
            width: 100%;
            height: 100%;
            object-fit: contain;
            display: block;
            transition: transform 0.6s cubic-bezier(0.4, 0, 0.2, 1);
        }

        .viz-box:hover img {
            transform: scale(1.02);
        }

        .lightbox {
            position: fixed;
            top: 0;
            left: 0;
            width: 100%;
            height: 100%;
            background: rgba(5, 8, 14, 0.95);
            backdrop-filter: blur(15px);
            z-index: 2000;
            display: none;
            justify-content: center;
            align-items: center;
            opacity: 0;
        }

        .lightbox-content {
            width: 90vw;
            height: 85vh;
            position: relative;
            display: flex;
            justify-content: center;
            align-items: center;
            box-shadow: 0 0 50px rgba(0, 0, 0, 0.5);
            border: 1px solid var(--glass-border);
            background: rgba(15, 23, 42, 0.8);
            border-radius: 12px;
            padding: 1rem;
        }

        .lightbox-img {
            width: 100%;
            height: 100%;
            display: block;
            object-fit: contain;
            cursor: pointer;
            transition: transform 0.3s;
        }

        .lightbox-img:hover {
            transform: scale(1.01);
        }

        .nav-controls {
            display: none;
        }

        .animate-up {
            opacity: 0;
            transform: translateY(30px);
            transition: opacity 0.8s ease-out, transform 0.8s ease-out;
        }
        
        /* Fallback: if JS is disabled or GSAP fails, show content */
        .no-js .animate-up,
        .js-fallback .animate-up {
            opacity: 1;
            transform: none;
        }

        @media print {
            .nav-controls,
            .bg-element {
                display: none;
            }

            body {
                overflow: visible;
            }

            .slide {
                height: 100vh;
                page-break-after: always;
                padding: 2rem;
                background: #fff !important;
                color: #000 !important;
            }

            .bg-card {
                background: none !important;
                color: #000 !important;
                border-color: #eee !important;
            }
        }
    
        .hud-status {
            position: fixed; top: 1.5rem; left: 3.5rem;
            display: flex; align-items: center; gap: 12px;
            font-family: 'Roboto Mono', monospace; font-size: 0.75rem; color: var(--accent); z-index: 1000;
        }
        .hud-dot {
            width: 6px; height: 6px; background-color: #ff00ea; border-radius: 50%;
            animation: blinkHUD 1.5s infinite;
        }
        @keyframes blinkHUD { 0%, 100% { opacity: 1; box-shadow: 0 0 8px #ff00ea; } 50% { opacity: 0.3; box-shadow: none; } }
        .hud-progress-container {
            position: fixed; top: 0; left: 0; width: 100%; height: 2px; background: rgba(255,255,255,0.05); z-index: 1001;
        }
        .hud-progress-bar {
            height: 100%; background: var(--accent); width: 0%; transition: width 0.5s ease; box-shadow: 0 0 10px var(--accent);
        }
        .hud-corners {
            position: fixed; top: 0; left: 0; width: 100%; height: 100%; pointer-events: none; z-index: 999;
        }
        .corner { position: absolute; width: 20px; height: 20px; border: 1px solid var(--accent); opacity: 0.5; }
        .corner.topleft { top: 20px; left: 20px; border-right: none; border-bottom: none; }
        .corner.topright { top: 20px; right: 20px; border-left: none; border-bottom: none; }
        .corner.bottomleft { bottom: 20px; left: 20px; border-right: none; border-top: none; }
        .corner.bottomright { bottom: 20px; right: 20px; border-left: none; border-top: none; }

        .presenter-card {
            box-shadow: inset 0 1px 0 rgba(255,255,255,0.1), inset 1px 0 0 rgba(255,255,255,0.05), 0 8px 32px rgba(0,0,0,0.3);
            position: relative; overflow: hidden;
        }
        .presenter-card::before {
            content: ''; position: absolute; top: 0; left: 0; right: 0; bottom: 0;
            background-image: url("data:image/svg+xml,%3Csvg viewBox='0 0 200 200' xmlns='http://www.w3.org/2000/svg'%3E%3Cfilter id='n'%3E%3CfeTurbulence type='fractalNoise' baseFrequency='0.85' numOctaves='3' stitchTiles='stitch'/%3E%3C/filter%3E%3Crect width='100%25' height='100%25' filter='url(%23n)' opacity='0.05'/%3E%3C/svg%3E");
            pointer-events: none; z-index: -1; opacity: 0.15;
        }

        .mag-nav {
            position: fixed; top: 80px; height: calc(100vh - 80px); width: 8vw; z-index: 1000;
            display: flex; align-items: center; justify-content: center;
            cursor: pointer; opacity: 0; transition: opacity 0.4s, background 0.4s;
            background: rgba(0, 242, 255, 0);
        }
        .mag-nav.left { left: 0; border-left: 2px solid transparent; }
        .mag-nav.right { right: 0; border-right: 2px solid transparent; }
        .mag-nav:hover { opacity: 1; background: rgba(0, 242, 255, 0.02); }
        .mag-nav.left:hover { border-left-color: var(--accent); }
        .mag-nav.right:hover { border-right-color: var(--accent); }
        .mag-nav span {
            color: var(--accent); font-family: 'Roboto Mono', monospace; font-size: 0.85rem;
            transform: rotate(-90deg); letter-spacing: 0.3em; white-space: nowrap; user-select: none;
        }
        .mag-nav.right span { transform: rotate(90deg); }

        .sidebar {
            position: fixed; top: 0; left: -380px; width: 380px; height: 100vh;
            background: rgba(15, 23, 42, 0.95); backdrop-filter: blur(25px);
            border-right: 2px solid var(--accent);
            z-index: 3000; transition: left 0.4s cubic-bezier(0.25, 1, 0.5, 1);
            display: flex; flex-direction: column; padding: 2rem 1.5rem;
            box-shadow: 20px 0 50px rgba(0,0,0,0.5);
        }
        .sidebar.open { left: 0; }
        .sidebar-close {
            position: absolute; top: 1.5rem; right: 1.5rem; color: var(--text-dim);
            cursor: pointer; transition: color 0.3s;
        }
        .sidebar-close:hover { color: var(--accent); }
        .sidebar-header {
            font-family: 'Outfit', sans-serif; font-weight: 700; color: var(--text-main);
            font-size: 1.2rem; margin-bottom: 2rem; padding-bottom: 1rem; border-bottom: 1px solid var(--glass-border);
        }
        .sidebar-list { display: flex; flex-direction: column; gap: 0.5rem; overflow-y: auto; }
        .sidebar-item {
            padding: 0.75rem; border-radius: 8px; cursor: pointer; transition: all 0.3s;
            display: flex; flex-direction: column; gap: 0.5rem;
            border: 1px solid transparent; background: rgba(255,255,255,0.02);
            position: relative;
        }
        .sidebar-item:hover { background: var(--accent-soft); border-color: var(--accent); }
        .sidebar-item.active { background: rgba(0, 242, 255, 0.15); border-left: 4px solid var(--accent); }
        .thumb-container {
            width: 100%; aspect-ratio: 16/9; background: var(--bg-deep);
            border-radius: 4px; overflow: hidden; position: relative;
            border: 1px solid rgba(255,255,255,0.1);
        }
        .sidebar-item-header { display: flex; justify-content: space-between; align-items: center; }
        .sidebar-item-num { font-family: 'Roboto Mono', monospace; font-size: 0.75rem; color: var(--accent); }
        .sidebar-item-title { font-size: 0.9rem; color: var(--text-main); line-height: 1.3; }

        .menu-trigger {
            cursor: pointer; display: flex; align-items: center; justify-content: center;
            color: var(--text-dim); transition: color 0.3s;
        }
        .menu-trigger:hover { color: var(--accent); }

        .summary-item, .roadmap-item {
            margin-bottom: 1.5rem;
            padding: 1.2rem;
            background: rgba(255,255,255,0.03);
            border-radius: 0.8rem;
            border: 1px solid rgba(255,255,255,0.05);
            transition: all 0.3s ease;
            display: flex;
            align-items: flex-start;
            gap: 1.2rem;
        }
        .summary-item:hover, .roadmap-item:hover {
            background: rgba(255,255,255,0.08);
            transform: translateX(10px);
            border-color: var(--accent);
        }
        .summary-item i, .roadmap-item i {
            color: var(--accent);
            flex-shrink: 0;
            margin-top: 4px;
            width: 1.6rem !important;
            height: 1.6rem !important;
        }

        .roadmap-item i {
            animation: pulse-rocket 2s infinite ease-in-out;
        }

        @keyframes pulse-rocket {
            0% { transform: scale(1); filter: drop-shadow(0 0 5px var(--accent)); }
            50% { transform: scale(1.1); filter: drop-shadow(0 0 15px var(--accent)); }
            100% { transform: scale(1); filter: drop-shadow(0 0 5px var(--accent)); }
        }
        .rocket-glow {
            position: absolute;
            width: 200px;
            height: 200px;
            background: radial-gradient(circle, var(--accent-soft) 0%, transparent 70%);
            z-index: -1;
        }
    </style>
</head>

<body>

    <div class="bg-element bg-1"></div>
    <div class="bg-element bg-2"></div>

    <div class="hud-status">
        <div class="menu-trigger" id="menuTrigger" title="Открыть оглавление"><i data-lucide="menu" style="width: 18px; height: 18px;"></i></div>
        <div class="hud-dot" style="margin-left: 8px;"></div>
    </div>
    <div class="hud-progress-container"><div class="hud-progress-bar" id="hudProgress"></div></div>
    <div class="hud-corners">
        <div class="corner topleft"></div><div class="corner topright"></div>
        <div class="corner bottomleft"></div><div class="corner bottomright"></div>
    </div>
    <div class="mag-nav left" id="magPrev"><span>// ПРЕДЫДУЩИЙ</span></div>
    <div class="mag-nav right" id="magNext"><span>// СЛЕДУЮЩИЙ</span></div>

    <aside class="sidebar" id="sidebar">
        <div class="sidebar-close" id="sidebarClose"><i data-lucide="x"></i></div>
        <div class="sidebar-header"><i data-lucide="layers" style="width: 20px; vertical-align: middle; margin-right: 8px;"></i>Оглавление</div>
        <div class="sidebar-list" id="sidebarList">
        </div>
    </aside>

    <div class="presentation" id="presentation">
    <!-- SLIDES_START -->"""


BASE_HTML_TAIL = """
    <!-- SLIDES_END -->
    </div>

    <div class="lightbox" id="lightbox">
        <div class="lightbox-content">
            <img src="" alt="Full view" class="lightbox-img" id="lightboxImg">
        </div>
    </div>

    <script>
        // Add JS detection class
        document.documentElement.classList.add('js');
        
        // Guard lucide
        if (typeof lucide !== 'undefined') {
            lucide.createIcons();
        }
        
        let currentSlide = 0;
        const presentation = document.getElementById('presentation');
        const totalSlides = presentation ? presentation.querySelectorAll('.slide').length : 0;

        function updateNavigation() {
            if (!presentation) return;
            presentation.style.transform = `translateX(-${currentSlide * 100}vw)`;
            document.querySelectorAll('.sidebar-item').forEach((item, index) => {
                if (index === currentSlide) item.classList.add('active');
                else item.classList.remove('active');
            });
            const progress = ((currentSlide + 1) / totalSlides) * 100;
            const hudBar = document.getElementById('hudProgress');
            if (hudBar) hudBar.style.width = progress + '%';
            const btnPrev = document.getElementById('magPrev');
            const btnNext = document.getElementById('magNext');
            if(btnPrev) btnPrev.style.display = currentSlide === 0 ? 'none' : 'flex';
            if(btnNext) btnNext.style.display = currentSlide === totalSlides - 1 ? 'none' : 'flex';

            const activeSlideItems = presentation.querySelectorAll('.slide')[currentSlide].querySelectorAll('.animate-up');
            if (activeSlideItems.length === 0) return;
            
            if (typeof gsap !== 'undefined') {
                gsap.fromTo(activeSlideItems,
                    { opacity: 0, y: 30 },
                    { opacity: 1, y: 0, duration: 0.8, stagger: 0.2, ease: "power2.out" }
                );
            } else {
                activeSlideItems.forEach(el => {
                    el.style.opacity = '1';
                    el.style.transform = 'none';
                });
                document.body.classList.add('js-fallback');
            }
        }

        // Fallback timeout
        setTimeout(() => {
            if (typeof gsap === 'undefined' && !document.body.classList.contains('js-fallback')) {
                document.querySelectorAll('.animate-up').forEach(el => {
                    el.style.opacity = '1';
                    el.style.transform = 'none';
                });
                document.body.classList.add('js-fallback');
                console.log('[Fallback] GSAP not loaded, showing content');
            }
        }, 2000);

        const mP = document.getElementById('magPrev');
        const mN = document.getElementById('magNext');
        if(mP) mP.addEventListener('click', () => { if (currentSlide > 0) { currentSlide--; updateNavigation(); } });
        if(mN) mN.addEventListener('click', () => { if (currentSlide < totalSlides - 1) { currentSlide++; updateNavigation(); } });

        const slides = presentation ? presentation.querySelectorAll('.slide') : [];
        const sidebarList = document.getElementById('sidebarList');
        slides.forEach((slide, index) => {
            let titleEl = slide.querySelector('.slide-title');
            let headingEl = slide.querySelector('.main-heading');
            let titleText = "Слайд " + (index + 1);
            if (titleEl && titleEl.innerText) {
                titleText = titleEl.innerText;
            } else if (headingEl && headingEl.innerText) {
                titleText = "Главный титул";
            }
            const item = document.createElement('div');
            item.className = 'sidebar-item';
            item.innerHTML = `
                <div class="sidebar-item-header">
                    <span class="sidebar-item-num">0${index+1} / 0${slides.length}</span>
                    <span class="sidebar-item-title">${titleText}</span>
                </div>
                <div class="thumb-container" id="thumb-container-${index}"></div>
            `;
            item.addEventListener('click', () => {
                currentSlide = index;
                updateNavigation();
            });
            sidebarList.appendChild(item);

            setTimeout(() => {
                const tContainer = document.getElementById(`thumb-container-${index}`);
                if (!tContainer) return;
                const clone = slide.cloneNode(true);
                clone.removeAttribute('id');
                clone.querySelectorAll('[id]').forEach(el => el.removeAttribute('id'));
                clone.style.width = '1920px';
                clone.style.height = '1080px';
                clone.style.position = 'absolute';
                clone.style.top = '0';
                clone.style.left = '0';
                clone.style.display = 'flex';
                clone.style.padding = '4rem 8rem';
                const scale = tContainer.clientWidth / 1920;
                clone.style.transform = `scale(${scale})`;
                clone.style.transformOrigin = 'top left';
                clone.style.pointerEvents = 'none';
                clone.querySelectorAll('.animate-up').forEach(el => {
                    el.style.opacity = '1';
                    el.style.transform = 'none';
                    el.style.clipPath = 'none';
                });
                tContainer.appendChild(clone);
            }, 100);
        });

        const sidebar = document.getElementById('sidebar');
        const menuTrigger = document.getElementById('menuTrigger');
        const sidebarClose = document.getElementById('sidebarClose');

        if (menuTrigger && sidebar) menuTrigger.addEventListener('click', () => { sidebar.classList.add('open'); });
        if (sidebarClose && sidebar) sidebarClose.addEventListener('click', () => { sidebar.classList.remove('open'); });

        // Call immediately since script is at end of body, DOM is ready
        updateNavigation();
        
        // Initialize Lucide icons when library is ready (deferred load)
        (function initLucide() {
            if (typeof lucide !== 'undefined') {
                lucide.createIcons();
            } else {
                setTimeout(initLucide, 100);
            }
        })();

        document.addEventListener('keydown', (e) => {
            if (e.key === 'ArrowRight' && currentSlide < totalSlides - 1) { currentSlide++; updateNavigation(); }
            if (e.key === 'ArrowLeft' && currentSlide > 0) { currentSlide--; updateNavigation(); }
            if (e.key === 'Escape') closeLightbox();
        });

        const lightbox = document.getElementById('lightbox');
        const lightboxImg = document.getElementById('lightboxImg');

        if (lightbox && lightboxImg) {
            document.querySelectorAll('.viz-box').forEach(box => {
                box.addEventListener('click', () => {
                    const img = box.querySelector('img');
                    if (img) {
                        lightboxImg.src = img.src;
                        lightbox.style.display = 'flex';
                        if (typeof gsap !== 'undefined') {
                            gsap.to(lightbox, { opacity: 1, duration: 0.4, ease: "power2.out" });
                            gsap.fromTo('.lightbox-content',
                                { scale: 0.8, y: 50 },
                                { scale: 1, y: 0, duration: 0.5, ease: "back.out(1.7)" }
                            );
                        } else {
                            lightbox.style.opacity = 1;
                        }
                    }
                });
            });

            function closeLightbox() {
                if (lightbox.style.display === 'flex') {
                    if (typeof gsap !== 'undefined') {
                        gsap.to(lightbox, {
                            opacity: 0,
                            duration: 0.3,
                            onComplete: () => { lightbox.style.display = 'none'; }
                        });
                    } else {
                        lightbox.style.display = 'none';
                    }
                }
            }

            lightboxImg.addEventListener('click', closeLightbox);
            lightbox.addEventListener('click', (e) => {
                if (e.target === lightbox) closeLightbox();
            });
        }

        /**
         * Автоматическое уплотнение контента при переполнении (Overflow Detection)
         */
        /**
         * Ультра-адаптивная функция плотности контента (Nuclear Fix)
         * Реализует многоходовое сжатие с запасом 10% и экстренным уменьшением шрифта.
         */
        async function applyDynamicDensity() {
            const panels = document.querySelectorAll('.analytical-panel');
            
            for (const panel of panels) {
                // 0. Сброс к дефолту
                panel.style.setProperty('--squeeze-factor', '1.0');
                panel.style.fontSize = ""; // сброс экстренного шрифта
                panel.style.overflowY = 'hidden';
                panel.style.paddingBottom = '0px';

                // Даем браузеру секунду на осознание (или requestAnimationFrame)
                await new Promise(r => requestAnimationFrame(r));
                
                let sH = panel.scrollHeight;
                let cH = panel.clientHeight;

                if (sH > cH + 1) {
                    // 1. Первая итерация: Сжатие интервалов с 7% запасом (relative safety margin)
                    // (cH/sH) * 0.93
                    let factor = (cH / sH) * 0.93;
                    factor = Math.max(0.60, factor); // Разрешаем сжимать до 60%
                    
                    panel.style.setProperty('--squeeze-factor', factor.toFixed(3));
                    panel.style.paddingBottom = '4px'; // Технический зазор
                    
                    // Ждем применения стилей
                    await new Promise(r => requestAnimationFrame(r));
                    
                    // 2. Если все еще не влезло (после MathJax и прочего)
                    if (panel.scrollHeight > panel.clientHeight + 1) {
                        // Пробуем "дожать" интервалы до абсолютного минимума
                        if (factor > 0.61) {
                            factor = Math.max(0.60, factor * 0.95);
                            panel.style.setProperty('--squeeze-factor', factor.toFixed(3));
                            await new Promise(r => requestAnimationFrame(r));
                        }
                        
                        // 3. ЭКСТРЕННЫЙ ПЛАН (Emergency Font Scaling)
                        // Если даже при факторе 0.60 не влезает - уменьшаем шрифт
                        if (panel.scrollHeight > panel.clientHeight + 1) {
                            panel.style.fontSize = "0.92em";
                            await new Promise(r => requestAnimationFrame(r));
                            
                            // Последняя попытка пересчитать фактор под новый размер шрифта
                            let finalSH = panel.scrollHeight;
                            let finalCH = panel.clientHeight;
                            if (finalSH > finalCH + 1) {
                                let finalFactor = Math.max(0.60, (finalCH / finalSH) * 0.95);
                                panel.style.setProperty('--squeeze-factor', finalFactor.toFixed(3));
                                await new Promise(r => requestAnimationFrame(r));
                            }
                        }
                    }

                    // 4. Финальный вердикт: если даже ядерный удар не помог - включаем скролл
                    if (panel.scrollHeight > panel.clientHeight + 2) {
                        panel.style.overflowY = 'auto';
                    }
                }
            }
        }

        window.addEventListener('load', () => {
            // Запуск через небольшую паузу для готовности MathJax
            setTimeout(applyDynamicDensity, 500);
        });
        
        let resizeTimer;
        window.addEventListener('resize', () => {
            clearTimeout(resizeTimer);
            resizeTimer = setTimeout(applyDynamicDensity, 200);
        });
    </script>
</body>
</html>
"""


def esc(s: str) -> str:
    """Экранирует строку для безопасного вставки в HTML."""
    return html.escape(str(s), quote=True)


def clean_text(text: str) -> str:
    """Нормализует текст: заменяет неразрывные пробелы, убирает лишние пробелы."""
    if not text:
        return ""
    # Заменяем неразрывный пробел на обычный
    text = text.replace("\xa0", " ")
    # Заменяем Wingdings bullets (U+F03E и подобные) на обычный bullet или удаляем
    text = re.sub(r"[\uf0b0-\uf0ff]", "", text)  # Private Use Area bullets
    # Collapse multiple spaces
    text = re.sub(r" +", " ", text)
    return text.strip()


class LLMCache:
    """Простой кэш для ответов LLM на основе файловой системы."""

    def __init__(self, cache_dir: str = None):
        self.cache_dir = Path(cache_dir) if cache_dir else BASE_DIR / ".cache" / "llm"
        self.cache_dir.mkdir(parents=True, exist_ok=True)

    def get(self, prompt: str, model: str) -> str | None:
        key = hashlib.sha256(f"{model}:{prompt}".encode()).hexdigest()
        path = self.cache_dir / f"{key}.json"
        if path.exists():
            try:
                with open(path, "r", encoding="utf-8") as f:
                    data = json.load(f)
                    age_days = (time.time() - data.get("ts", 0)) / 86400
                    if age_days < 30:
                        return data.get("response")
            except Exception as e:
                logger.warning(f"Ошибка чтения кэша LLM: {e}")
        return None

    def set(self, prompt: str, model: str, response: str):
        key = hashlib.sha256(f"{model}:{prompt}".encode()).hexdigest()
        path = self.cache_dir / f"{key}.json"
        try:
            with open(path, "w", encoding="utf-8") as f:
                json.dump(
                    {
                        "prompt": prompt,
                        "model": model,
                        "response": response,
                        "ts": time.time(),
                    },
                    f,
                    ensure_ascii=False,
                    indent=2,
                )
        except Exception as e:
            logger.error(f"Ошибка сохранения кэша LLM: {e}")


# ==============================================================================
# БЛОК КОНСТАНТ ОФОРМЛЕНИЯ (DESIGN CONSTANTS)
# ==============================================================================
DESIGN_CONFIG = {
    "icon_size": "1.6rem",  # Relative size for regular list item icons
    "bullet_lists": {
        "enabled": True,
        "icon_map": {
            "•": "chevron-right",
            "◦": "chevron-right",
            "▪": "square",
            "-": "chevron-right",
            "–": "chevron-right",
            "—": "chevron-right",
            "*": "star",
            "": "chevron-right",
            "·": "dot",
            "": "chevron-right",
            "ü": "check",
        },
        "icon_size": "1.6rem",
        "indent": "2rem",
        "border_left": "2px solid rgba(0, 242, 255, 0.15)",
        "background": "rgba(0, 242, 255, 0.03)",
    },
    "grid": {
        "max_columns": 4, # Максимальное число колонок в сетке изображений
        "col_span_threshold": 1.2,  # aspect >= 1.2 → занимает 2 колонки
        "row_span_threshold": 0.833, # aspect <= 0.833 → занимает 2 строки
        "fallback_min_col_width": "200px",
    },
    "caption_search": {
        "vertical_range_mm": 25, # Диапазон поиска подписей вокруг картинки (мм)
        "max_gap_mm": 20,         # Максимальный зазор до подписи
        "horizontal_overlap_ratio": 0.6,
        "overlap_tolerance": 0.4,
        "priority": "above",     # Приоритетное расположение подписи: сверху
    },
    "formula": {
        "mathjax_path": "libs/mathjax/tex-mml-svg.js", # Путь к локальному MathJax (SVG версия)
        "fallback_font": "Roboto Mono",
        "fallback_font_size": "0.9em",
        "padding": "1.5rem",
        "background": "rgba(255, 255, 255, 0.01)",
        "border_color": "rgba(0, 242, 255, 0.15)",
    },
    "layout": {
        "caption_height_px": 50,  # Резерв высоты под заголовок рисунка внутри Grid-ячейки
        "text_panel_ratio": 0.35, # Стандартная пропорция ширины текстовой панели
    },
    "paths": {
        "logo_white": "logo/white.png",
        "media_output": "media/smart_present",
        "media_output_full": str(BASE_DIR / "web_demo" / "media" / "smart_present"),
    },
    "icon_mapping": {
        # Маппинг ключевых слов на иконки Lucide
        "activity": ["динамика", "поле", "процесс", "геодинамика"],
        "droplet": ["нефть", "жидкость", "поток", "вода"],
        "layers": ["стратиграфия", "пласт", "разрез", "толща", "литология", "фондоформ"],
        "bar-chart": ["результат", "статистика", "данные", "анализ", "экономика", "итог"],
        "compass": ["направление", "азимут", "ориентация", "σhmax", "нmax", "тренд"],
        "target": ["цел", "перспектив", "направлен"],
        "list-todo": ["задач", "план", "постановк", "задан", "roadmap"],
        "database": ["модел", "сетк", "данн", "3d", "ячеек"],
        "cpu": ["автоматизац", "алгоритм", "расчет", "abaqus", "внедрен"],
        "alert-triangle": ["риск", "проблем", "опасност", "вниман", "предупрежден"],
        "refresh-ccw": ["ппд", "эффективност", "обработк"],
        "sliders": ["оптимизац", "параметр", "настройк"],
        "maximize": ["разм", "диаметр", "толщ", "глубин", "высот", "длин", "ширин"],
        "map-pin": ["регион", "месторожден", "район", "западн", "сибир", "участ"],
        "tower-control": ["скважин", "скв", "забой", "усть", "ствол"],
        "test-tube-2": ["испытан", "образц", "эксперимент", "лаборат"],
        "wrench": ["установк", "инструмент", "аппарат", "датчик", "прибор"],
        "zap": ["чувствительн", "влиян", "отклик", "эффект", "фактор"],
        "git-merge": ["нормирова", "приведен", "коррекц", "сопоставл"],
        "file-check": ["отчет", "регламент", "утвержден", "формат"],
        "info": ["информ", "инфо", "описан", "сведен", "справоч", "примечан"],
        "box": ["модель", "коробка", "box"],
    },
    "ai_config": {
        "env_path": str(BASE_DIR / ".env"),
        "default_model": "google/gemini-flash-1.5",
    },
    "STATIC_ASSETS_EMBED": True,
}


class PPTConverter:
    """Конвертер PowerPoint-презентаций в HTML с современным дизайном.

    Attributes:
        ppt_path: Путь к исходному .pptx файлу.
        output_html: Путь для выходного HTML.
        slides_data: Список данных о слайдах.
        stats: Статистика конвертации.
    """

    def __init__(self, ppt_path: str, output_html: str = None):
        self.ppt_path = ppt_path
        self.template_path = None  # Не используется, шаблон вшит в код
        self.output_html = output_html
        self.slides_data = []
        self.stats = {
            "total_slides": 0,
            "images_ok": 0,
            "images_fail": 0,
            "tables": 0,
            "formulas": 0,
            "ole_skipped": 0,
        }

    def _get_data_uri(self, file_path: Path) -> str:
        """Превращает файл в Base64 Data URI."""
        if not file_path.exists():
            logger.warning(f"Файл не найден для эмбеддинга: {file_path}")
            return ""
        
        mime_type, _ = mimetypes.guess_type(str(file_path))
        if not mime_type:
            mime_type = "application/octet-stream"
            
        with open(file_path, "rb") as f:
            data = f.read()
            b64 = base64.b64encode(data).decode('utf-8')
            return f"data:{mime_type};base64,{b64}"

    def _get_file_content(self, file_path: Path) -> str:
        """Читает текстовый файл."""
        if not file_path.exists():
            logger.warning(f"Текстовый файл не найден: {file_path}")
            return ""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()

    def _inline_css_fonts(self, css_path: Path) -> str:
        """Читает CSS и вшивает шрифты в него через Data URI."""
        content = self._get_file_content(css_path)
        if not content:
            return ""
        
        # Регулярка для поиска url(...)
        def fill_url(match):
            url_path = match.group(1).strip("'\"")
            # Пути в fonts.css относительны самой таблицы стилей
            full_path = (css_path.parent / url_path).resolve()
            data_uri = self._get_data_uri(full_path)
            # Match browser expectation for format(...)
            return f"url('{data_uri}')"
        
        return re.sub(r"url\((.*?)\)", fill_url, content)

    def _spatial_sort(self, items: list, threshold_mm: float = 20.0) -> list:
        """
        Сортировка изображений с использованием 'жадных строк' (greedy rows).
        Для элементов с pos = (left, top, width, height).
        """
        if not items:
            return []
        
        threshold = threshold_mm * 36000 # 1 mm ~ 36000 EMU
        
        # Сортируем сначала по Y (top)
        sorted_by_y = sorted(items, key=lambda x: x.get("pos", (0, 0, 0, 0))[1])
        
        rows = []
        if sorted_by_y:
            current_row = [sorted_by_y[0]]
            last_y = sorted_by_y[0].get("pos", (0, 0, 0, 0))[1]
            
            for item in sorted_by_y[1:]:
                curr_y = item.get("pos", (0, 0, 0, 0))[1]
                if abs(curr_y - last_y) < threshold:
                    current_row.append(item)
                else:
                    rows.append(sorted(current_row, key=lambda x: x.get("pos", (0, 0, 0, 0))[0]))
                    current_row = [item]
                    last_y = curr_y
            rows.append(sorted(current_row, key=lambda x: x.get("pos", (0, 0, 0, 0))[0]))
            
        final_sorted = []
        for row in rows:
            final_sorted.extend(row)
        return final_sorted

    def _spatial_sort_strict(self, items: list) -> list:
        """
        Строгая вертикальная сортировка элементов контента.
        Используется для соблюдения хронологического порядка (сверху вниз)
        в текстовой панели (текст, таблицы, формулы).
        """
        if not items:
            return []
        # Сортируем по координате Y (top), затем по X (left)
        return sorted(items, key=lambda x: (x.get("pos", (0, 0))[0], x.get("pos", (0, 0))[1]))

    def get_icon_by_text(self, text: str) -> str:
        """Возвращает идентификатор иконки на основе текста (по ключевым словам)."""
        text = text.lower()
        for icon, keywords in DESIGN_CONFIG["icon_mapping"].items():
            if any(kw in text for kw in keywords):
                return icon
        return "chevron-right"

    def load_env(self) -> dict:
        """Загружает переменные окружения из .env файла."""
        env = {}
        path = DESIGN_CONFIG["ai_config"]["env_path"]
        if os.path.exists(path):
            with open(path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    if "=" in line and not line.startswith("#"):
                        k, v = line.split("=", 1)
                        env[k.strip()] = v.strip().strip('"').strip("'")
        return env

    def call_ai(self, prompt: str) -> str | None:
        """Вызывает LLM через OpenRouter API с кэшированием."""
        env = self.load_env()
        api_key = env.get("OPENROUTER_API_KEY")
        model = env.get("MODEL", DESIGN_CONFIG["ai_config"]["default_model"])

        if not api_key:
            logger.warning("OPENROUTER_API_KEY не найден в .env")
            return None

        cache = LLMCache()
        cached = cache.get(prompt, model)
        if cached:
            logger.info(f"[LLM] Ответ получен из кэша для prompt: {prompt[:50]}...")
            return cached

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }
        payload = {"model": model, "messages": [{"role": "user", "content": prompt}]}

        try:
            response = requests.post(
                "https://openrouter.ai/api/v1/chat/completions",
                headers=headers,
                data=json.dumps(payload),
                timeout=30,
            )
            if response.status_code == 200:
                result = response.json()
                content = result["choices"][0]["message"]["content"]
                cache.set(prompt, model, content)
                return content
            else:
                logger.error(
                    f"Ошибка API: {response.status_code} - {response.text[:200]}"
                )
        except Exception as e:
            logger.error(f"Исключение при вызове ИИ: {e}")
        return None

    def process_txt_files(self) -> None:
        """Обрабатывает вспомогательные текстовые файлы с использованием AI для структурирования."""
        files = {
            "Выводы.txt": {
                "type": "conclusions",
                "prompt": "Проанализируй текст и извлеки ключевые факты. СТРОГАЯ СТРУКТУРА: 1. Выполненные работы (МАКСИМУМ 6 пунктов): перечисли основные этапы и действия. 2. Результаты (МАКСИМУМ 6 пунктов): перечисли конкретные выводы, достижения и показатели. ПРАВИЛА: - СТРОГО соблюдай хронологический порядок событий. - Текст должен быть максимально сжатым, строгим и технически точным. - Объединяй весь опыт, не пропуская значимых деталей (даты, цифры), но формулируй их крайне емко. - Удали вводные слова, пояснения и 'воду'. ФОРМАТ: Верни только JSON-массив из кратких тезисов (сначала до 6 пунктов по работам, затем до 6 пунктов по результатам). Текст: {content}",
            },
            "Направление дальнейших исследований.txt": {
                "type": "research",
                "prompt": "Извлеки основные направления дальнейших исследований. ПОРЯДОК: Сначала планируемые действия, затем ожидаемые эффекты и цели. ПРАВИЛА: - Соблюдай логическую и хронологическую последовательность. - Максимально краткий и сухой научный стиль. - Сохрани всю фактологическую базу данных. ФОРМАТ: Верни только JSON-массив строк. Текст: {content}",
            },
        }

        for filename, config in files.items():
            if not os.path.exists(filename):
                continue
            logger.info(f"Обработка текстового файла: {filename}")
            with open(filename, "r", encoding="utf-8") as f:
                content = f.read()

            items = []
            ai_used = False
            api_key = self.load_env().get("OPENROUTER_API_KEY")

            if api_key:
                prompt = config["prompt"].format(content=content[:8000])
                ai_response = self.call_ai(prompt)
                if ai_response:
                    try:
                        # Clean markdown code fences if present
                        cleaned = ai_response.strip()
                        if cleaned.startswith("```json"):
                            cleaned = cleaned.split("```json", 1)[1]
                        if cleaned.startswith("```"):
                            cleaned = cleaned.split("```", 1)[1]
                        if "```" in cleaned:
                            cleaned = cleaned.split("```", 1)[0]
                        cleaned = cleaned.strip()

                        parsed = json.loads(cleaned)
                        if isinstance(parsed, list) and all(
                            isinstance(x, str) for x in parsed
                        ):
                            items = [p.strip() for p in parsed if p.strip()]
                            ai_used = True
                            logger.info(
                                f"[AI] Получено {len(items)} пунктов из {filename}"
                            )
                        else:
                            logger.warning(
                                f"[AI] Ответ не является массивом строк, используем fallback"
                            )
                    except json.JSONDecodeError as e:
                        logger.warning(
                            f"[AI] Не удалось распарсить JSON: {e}, используем fallback"
                        )

            if not items:
                # Fallback: разбиваем по строкам
                paragraphs = [p.strip() for p in content.split("\n") if p.strip()]
                items = paragraphs
                if not ai_used:
                    logger.info(
                        f"[Fallback] Использовано разбиение по строкам для {filename} ({len(items)} пунктов)"
                    )

            if config["type"] == "conclusions":
                processed_items = []
                for item in items:
                    clean_item = clean_text(item)
                    icon = self.get_icon_by_text(clean_item)
                    processed_items.append(
                        f"<div class='summary-item'><i data-lucide='{icon}'></i> <div class='list-text'>{esc(clean_item)}</div></div>"
                    )

                mid = (len(processed_items) + 1) // 2
                left_html = "".join(processed_items[:mid])
                right_html = "".join(processed_items[mid:])

                self.slides_data.append(
                    {
                        "title": "Выводы и результаты этапа",
                        "layout_type": "conclusions_dual",
                        "left_html": left_html,
                        "right_html": right_html,
                        "content_items": [],
                    }
                )
            elif config["type"] == "research":
                items_html = []
                for item in items:
                    clean_item = clean_text(item)
                    items_html.append(
                        f"<div class='roadmap-item'><i data-lucide='rocket'></i> <div class='list-text'>{esc(clean_item)}</div></div>"
                    )
                self.slides_data.append(
                    {
                        "title": "Направление дальнейших исследований",
                        "layout_type": "research_roadmap",
                        "content_html": "".join(items_html),
                        "content_items": [],
                    }
                )

    def _save_image_with_white_bg(self, img_blob: bytes, output_path: str) -> bool:
        """Сохраняет изображение, заменяя прозрачный фон на белый. Возвращает True при успехе."""
        try:
            if not img_blob or len(img_blob) == 0:
                logger.warning(f"Пустой blob изображения, пропускаем {output_path}")
                return False

            img = Image.open(io.BytesIO(img_blob)).convert("RGBA")
            white_bg = Image.new("RGBA", img.size, (255, 255, 255, 255))
            composite = Image.alpha_composite(white_bg, img)
            composite = composite.convert("RGB")
            composite.save(output_path, "PNG")

            # Проверяем, что файл создан и не пустой
            if os.path.getsize(output_path) == 0:
                logger.warning(f"Файл {output_path} имеет размер 0 байт, удаляем")
                try:
                    os.remove(output_path)
                except Exception:
                    pass
                return False

            return True
        except Exception as e:
            logger.error(f"Ошибка сохранения {output_path}: {e}")
            # Удаляем частично созданный файл, если он есть
            if os.path.exists(output_path):
                try:
                    os.remove(output_path)
                except Exception:
                    pass
            return False

    def get_best_layout(self, visuals: list, container_w: float, container_h: float):
        """
        Минимизирует пустое пространство для 1-5 рисунков.
        Для 6-8 рисунков использует фиксированную сетку.
        Возвращает: (rows, cols, grid_styles)
        """
        n = len(visuals)
        if n == 0: return 1, 1, [], "1fr", "1fr"
        
        caption_h = DESIGN_CONFIG["layout"]["caption_height_px"]
        
        # 6-8 рисунков: строгая сетка в 2 строки
        if n >= 6:
            cols = math.ceil(n / 2)
            grid_styles = []
            for i in range(n):
                r, c = i // cols, i % cols
                grid_styles.append(f"grid-row: {r+1} / span 1; grid-column: {c+1} / span 1;")
            col_tmpl = " ".join(["1fr"] * cols)
            return 2, cols, grid_styles, "1fr 1fr", col_tmpl

        aspects = []
        for v in visuals:
            _, _, vw, vh = v["pos"]
            asp = vw / vh if vh > 0 else 1.33
            aspects.append(asp)

        # Шаблоны для 1-5 рисунков: (rows, cols, items_layout)
        templates_n = {
            1: [(1, 1, [{'r':0, 'c':0, 'rs':1, 'cs':1}])],
            2: [
                (1, 2, [{'r':0, 'c':0, 'rs':1, 'cs':1}, {'r':0, 'c':1, 'rs':1, 'cs':1}]), # Колонки
                (2, 1, [{'r':0, 'c':0, 'rs':1, 'cs':1}, {'r':1, 'c':0, 'rs':1, 'cs':1}]), # Строки
            ],
            3: [
                (1, 3, [{'r':0, 'c':0, 'rs':1, 'cs':1}, {'r':0, 'c':1, 'rs':1, 'cs':1}, {'r':0, 'c':2, 'rs':1, 'cs':1}]), # 3 колонки
                (2, 2, [{'r':0, 'c':0, 'rs':2, 'cs':1}, {'r':0, 'c':1, 'rs':1, 'cs':1}, {'r':1, 'c':1, 'rs':1, 'cs':1}]), # 1 высокая + 2 справа
                (2, 2, [{'r':0, 'c':1, 'rs':2, 'cs':1}, {'r':0, 'c':0, 'rs':1, 'cs':1}, {'r':1, 'c':0, 'rs':1, 'cs':1}]), # 1 высокая справа
                (2, 2, [{'r':0, 'c':0, 'rs':1, 'cs':2}, {'r':1, 'c':0, 'rs':1, 'cs':1}, {'r':1, 'c':1, 'rs':1, 'cs':1}]), # 1 широкая сверху
            ],
            4: [
                (2, 2, [{'r':0, 'c':0, 'rs':1, 'cs':1}, {'r':0, 'c':1, 'rs':1, 'cs':1}, {'r':1, 'c':0, 'rs':1, 'cs':1}, {'r':1, 'c':1, 'rs':1, 'cs':1}]), # 2x2
                (1, 4, [{'r':0, 'c':0, 'rs':1, 'cs':1}, {'r':0, 'c':1, 'rs':1, 'cs':1}, {'r':0, 'c':2, 'rs':1, 'cs':1}, {'r':0, 'c':3, 'rs':1, 'cs':1}]), # 4 колонки
            ],
            5: [
                (2, 3, [{'r':0, 'c':0, 'rs':2, 'cs':1}, {'r':0, 'c':1, 'rs':1, 'cs':1}, {'r':0, 'c':2, 'rs':1, 'cs':1}, {'r':1, 'c':1, 'rs':1, 'cs':1}, {'r':1, 'c':2, 'rs':1, 'cs':1}]), # 2x3 с 1 высоким
                (3, 2, [{'r':0, 'c':0, 'rs':1, 'cs':2}, {'r':1, 'c':0, 'rs':1, 'cs':1}, {'r':1, 'c':1, 'rs':1, 'cs':1}, {'r':2, 'c':0, 'rs':1, 'cs':1}, {'r':2, 'c':1, 'rs':1, 'cs':1}]), # 3x2 с 1 широким
            ]
        }
        
        candidates = templates_n.get(n, [])
        best_score = -1
        best_layout = (1, n, [{'r':0, 'c':i, 'rs':1, 'cs':1} for i in range(n)]) # fallback
        
        for rows, cols, items in candidates:
            cell_w = container_w / cols
            cell_h = container_h / rows
            current_total_area = 0
            
            for i, item in enumerate(items):
                cw = cell_w * item['cs']
                ch = cell_h * item['rs']
                img_ch = ch - caption_h
                if img_ch <= 0: continue
                
                asp = aspects[i]
                scale = min(cw / asp, img_ch)
                current_total_area += (scale * asp) * scale
            
            if current_total_area > best_score:
                best_score = current_total_area
                best_layout = (rows, cols, items)
        
        # Calculate weighted grid templates
        r_res, c_res, i_res = best_layout
        grid_styles = []
        for item in i_res:
            style = f"grid-row: {item['r']+1} / span {item['rs']}; grid-column: {item['c']+1} / span {item['cs']};"
            grid_styles.append(style)
            
        # Heuristic for column/row weights
        col_weights = ["1fr"] * c_res
        row_weights = ["1fr"] * r_res
        
        if n == 2 and c_res == 2:
            a1, a2 = aspects[0], aspects[1]
            if a1 < 0.7 or a2 < 0.7:
                w1 = max(0.6, min(1.4, a1))
                w2 = max(0.6, min(1.4, a2))
                col_weights = [f"{w1:.1f}fr", f"{w2:.1f}fr"]
        elif n == 3 and c_res == 3:
            w = [max(0.7, min(1.3, a)) for a in aspects]
            col_weights = [f"{v:.1f}fr" for v in w]
            
        return r_res, c_res, grid_styles, " ".join(row_weights), " ".join(col_weights)

    def _is_slide_number(self, text: str) -> bool:
        """Проверяет, является ли текст номером слайда (например, "3" или "3/10")."""
        if re.match(r"^\d+$", text.strip()):
            return True
        if re.match(r"^\d+\s*/\s*\d+$", text.strip()):
            return True
        return False

    def _split_text_into_items(self, text: str) -> list:
        """Разбирает многострочный текст на отдельные пункты.
        Возвращает список dict: [{"text": str, "is_bullet": bool, "bullet_char": str|None, "level": int}]
        Префиксы списков (цифры, bullets, дефисы) удаляются из текста, но тип маркера запоминается.
        """
        if not text:
            return []
        lines = text.split("\n")
        items = []
        # Символы маркера из конфига (ключи словаря)
        bullet_chars = set(
            DESIGN_CONFIG.get("bullet_lists", {}).get("icon_map", {}).keys()
        )
        # Расширенный паттерн для удаления префиксов: цифры с точками, bullets, дефисы
        # Исключаем буквы из длинных слов ([a-zA-Zа-яА-Я]+), оставляя только одиночные буквы-маркеры (a., B., а.)
        prefix_pattern = r"^\s*(\d+[\.\)]\s*|[a-zа-я][\.\)]\s*|[•◦▪\-–—*·ü]\s*)+\s*"
        for line in lines:
            stripped = line.strip()
            if not stripped:
                continue
            first_char = stripped[0]
            # КОРРЕКЦИЯ: Более гибкое определение списков (любая нумерация или пунктировка)
            is_bullet = (
                (first_char in bullet_chars) or 
                bool(re.match(r"^\s*(\d+[\.\)]|[a-zA-Zа-яА-Я][\.\)]|\([\d\w]\))\s*", stripped)) or
                bool(re.match(r"^\s*[IVXLCDMivxlcdm]+\.\s*", stripped))
            )
            
            cleaned = re.sub(prefix_pattern, "", stripped).strip()
            if cleaned:
                items.append(
                    {
                        "text": cleaned,
                        "is_bullet": is_bullet,
                        "bullet_char": first_char if is_bullet else None,
                        "level": 0,
                    }
                )
        return items

    def _extract_captions_from_shapes(self, slide, slide_info: dict) -> None:
        """Извлекает подписи к изображениям из текстовых блоков.
        Поиск сначала выполняется над изображением (priority=above), если не найдено —
        выполняется fallback поиск под изображением.
        Выбранный в качестве подписи текст исключается из общего списка plain_text.
        """
        cfg = DESIGN_CONFIG.get("caption_search", {})
        vert_range_mm = cfg.get("vertical_range_mm", 15)
        max_gap_mm = cfg.get("max_gap_mm", 10)
        horiz_ratio = cfg.get("horizontal_overlap_ratio", 0.4)
        overlap_tol = cfg.get("overlap_tolerance", 0.5)
        priority = cfg.get("priority", "above")
        vert_range_emu = int(vert_range_mm * 36000)
        max_gap_emu = int(max_gap_mm * 36000)

        slide_num = slide_info.get("slide_num", "?")
        text_shapes_with_pos = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                txt = shape.text.strip()
                if txt and not self._is_slide_number(txt) and len(txt) < 250:
                    text_shapes_with_pos.append(
                        {
                            "text": txt,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height,
                        }
                    )

        used_texts = set()
        for vis in slide_info["visuals"]:
            v_left, v_top, v_w, v_h = vis["pos"]
            v_center_x = v_left + v_w / 2
            best_match = None
            best_score = 0.0
            second_best_score = 0.0

            # --- Pass 1: search ABOVE image ---
            search_top = v_top - vert_range_emu

            for ts in text_shapes_with_pos:
                if ts["text"] in used_texts:
                    continue
                ts_top = ts["top"]
                ts_bottom = ts["top"] + ts["height"]
                ts_center_x = ts["left"] + ts["width"] / 2
                ts_h = ts["height"]

                # Too high
                if ts_bottom < search_top:
                    continue

                v_gap = v_top - ts_bottom
                abs_gap = abs(v_gap)
                if abs_gap > max_gap_emu:
                    continue

                # Overlap handling
                if v_gap < 0:
                    if ts_h > 0:
                        overlap_above = (v_top - ts_top) / ts_h
                        if overlap_above < overlap_tol:
                            continue
                    else:
                        continue

                # Horizontal matching
                h_offset = abs(ts_center_x - v_center_x)
                if h_offset > v_w * horiz_ratio:
                    continue
                # h_score: denominator fixed 0.5 per plan
                h_thresh_score = v_w * 0.5 if v_w > 0 else 1
                h_score = max(0.0, 1.0 - (h_offset / h_thresh_score))

                # v_score
                v_score = (
                    max(0.0, 1.0 - (abs_gap / max_gap_emu)) if max_gap_emu > 0 else 0.0
                )

                total_score = 0.7 * v_score + 0.3 * h_score

                if total_score > best_score:
                    second_best_score = best_score
                    best_score = total_score
                    best_match = ts
                elif total_score > second_best_score:
                    second_best_score = total_score

            # --- Pass 2: fallback BELOW image ---
            if priority == "above" and (best_match is None or best_score <= 0.6):
                v_bottom = v_top + v_h
                for ts in text_shapes_with_pos:
                    if ts["text"] in used_texts:
                        continue
                    ts_top = ts["top"]
                    ts_center_x = ts["left"] + ts["width"] / 2

                    # Horizontal check
                    h_offset = abs(ts_center_x - v_center_x)
                    if h_offset > v_w * horiz_ratio:
                        continue

                    # Text must start at or below image bottom, within max_gap
                    v_gap_below = ts_top - v_bottom
                    if v_gap_below < 0 or v_gap_below > max_gap_emu:
                        continue

                    # Scores
                    v_score = (
                        max(0.0, 1.0 - (v_gap_below / max_gap_emu))
                        if max_gap_emu > 0
                        else 0.0
                    )
                    h_thresh_score = v_w * 0.5 if v_w > 0 else 1
                    h_score = max(0.0, 1.0 - (h_offset / h_thresh_score))
                    total_score = 0.7 * v_score + 0.3 * h_score

                    if total_score > best_score:
                        second_best_score = best_score
                        best_score = total_score
                        best_match = ts
                    elif total_score > second_best_score:
                        second_best_score = total_score

            if best_match and best_score > 0.6:
                vis["caption"] = best_match["text"]
                used_texts.add(best_match["text"])
                
                # КОРРЕКЦИЯ: Надежная фильтрация через нормализацию пробелов
                def normalize(s: str) -> str:
                    return " ".join(s.split()).strip().lower()
                
                target = normalize(best_match["text"])
                
                new_items = []
                for item in slide_info.get("content_items", []):
                    if item["type"] == "text":
                        new_paras = []
                        for p in item["data"]:
                            p_text = normalize("".join(s for s in p if isinstance(s, str) and not s.startswith("[[[")))
                            if p_text != target:
                                new_paras.append(p)
                        if new_paras:
                            item["data"] = new_paras
                            new_items.append(item)
                    else:
                        new_items.append(item)
                slide_info["content_items"] = new_items
                if best_score - second_best_score < 0.1 and second_best_score > 0:
                    logger.warning(
                        f"Slide {slide_num}: ambiguous caption (score diff {best_score - second_best_score:.3f})"
                    )
            else:
                logger.warning(
                    f"Slide {slide_num}: caption not found for visual at (x={v_left}, y={v_top})"
                )

    def _omml_to_mathml(self, omath_elem) -> str:
        """Рекурсивно конвертирует OMML в MathML (без внешних библиотек)."""
        OMML_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/math}"
        
        def token_type(txt: str) -> str:
            if txt.isdigit() or txt.replace('.', '').isdigit():
                return "mn"
            ops = set("+-*/=<>≤≥≈≡∑∏∫∂∇±×÷∈∉⊂⊃∪∩∧∨¬∞∂")
            if txt in ops:
                return "mo"
            return "mi"

        def convert(elem):
            if elem is None:
                return ""
            tag = elem.tag.replace(OMML_NS, '')
            if tag == 'oMath' or tag == 'oMathPara':
                children = "".join(convert(c) for c in elem)
                return f'<math xmlns="http://www.w3.org/1998/Math/MathML">{children}</math>'
            elif tag == 'r':
                texts = []
                for t in elem.findall(f'.//{OMML_NS}t'):
                    if t.text:
                        texts.append(t.text)
                txt = ''.join(texts)
                if txt:
                    tt = token_type(txt)
                    return f'<{tt}>{html.escape(txt)}</{tt}>'
                return ''
            elif tag == 't':
                return html.escape(elem.text or '')
            elif tag == 'sSub':
                base = convert(elem.find(f'{OMML_NS}e'))
                sub = convert(elem.find(f'{OMML_NS}sub'))
                return f'<msub><mrow>{base}</mrow><mrow>{sub}</mrow></msub>'
            elif tag == 'sSup':
                base = convert(elem.find(f'{OMML_NS}e'))
                sup = convert(elem.find(f'{OMML_NS}sup'))
                return f'<msup><mrow>{base}</mrow><mrow>{sup}</mrow></msup>'
            elif tag == 'sSubSup':
                base = convert(elem.find(f'{OMML_NS}e'))
                sub = convert(elem.find(f'{OMML_NS}sub'))
                sup = convert(elem.find(f'{OMML_NS}sup'))
                return f'<msubsup><mrow>{base}</mrow><mrow>{sub}</mrow><mrow>{sup}</mrow></msubsup>'
            elif tag == 'f':
                num = convert(elem.find(f'{OMML_NS}num'))
                den = convert(elem.find(f'{OMML_NS}den'))
                return f'<mfrac><mrow>{num}</mrow><mrow>{den}</mrow></mfrac>'
            elif tag == 'rad':
                deg = elem.find(f'{OMML_NS}deg')
                e = convert(elem.find(f'{OMML_NS}e'))
                if deg is not None:
                    deg_val = convert(deg)
                    return f'<mroot><mrow>{e}</mrow><mrow>{deg_val}</mrow></mroot>'
                else:
                    return f'<msqrt><mrow>{e}</mrow></msqrt>'
            elif tag == 'd':
                # Delimiters (parentheses, brackets)
                dPr = elem.find(f'{OMML_NS}dPr')
                beg = '('
                end = ')'
                if dPr is not None:
                    begChr = dPr.find(f'{OMML_NS}begChr')
                    endChr = dPr.find(f'{OMML_NS}endChr')
                    if begChr is not None and begChr.get('val'):
                        beg = begChr.get('val')
                    if endChr is not None and endChr.get('val'):
                        end = endChr.get('val')
                content = convert(elem.find(f'{OMML_NS}e'))
                return f'<mrow><mo>{html.escape(beg)}</mo><mrow>{content}</mrow><mo>{html.escape(end)}</mo></mrow>'
            elif tag == 'eqArr':
                # Система уравнений
                rows = []
                for e in elem.findall(f'{OMML_NS}e'):
                    rows.append(convert(e))
                return f'<mtable>{"".join(f"<mtr><mtd><mrow>{r}</mrow></mtd></mtr>" for r in rows)}</mtable>'
            elif tag == 'nary':
                # Sum, integral, etc. (munderover)
                chr_elem = elem.find(f'{OMML_NS}chr')
                op = chr_elem.get('val') if chr_elem is not None else '∑'
                sub = convert(elem.find(f'{OMML_NS}sub'))
                sup = convert(elem.find(f'{OMML_NS}sup'))
                e = convert(elem.find(f'{OMML_NS}e'))
                if sub or sup:
                    # munderover expects 3 arguments: base, lower, upper
                    # We ensure each is a single child using <mrow>
                    s = sub if sub else '<mrow></mrow>'
                    sp = sup if sup else '<mrow></mrow>'
                    return f'<munderover><mrow><mo>{html.escape(op)}</mo></mrow><mrow>{s}</mrow><mrow>{sp}</mrow></munderover><mrow>{e}</mrow>'
                else:
                    return f'<mo>{html.escape(op)}</mo><mrow>{e}</mrow>'
            elif tag == 'acc':
                # Accent (hat, bar) - mover
                chr_elem = elem.find(f'{OMML_NS}chr')
                acc = chr_elem.get('val') if chr_elem is not None else '̂'
                base = convert(elem.find(f'{OMML_NS}e'))
                return f'<mover><mrow>{base}</mrow><mrow><mo>{html.escape(acc)}</mo></mrow></mover>'
            elif tag == 'box':
                return f'<mrow>{convert(elem.find(f"{OMML_NS}e"))}</mrow>'
            elif tag == 'e':
                return "".join(convert(c) for c in elem)
            else:
                # Рекурсивно обрабатываем дочерние элементы
                return "".join(convert(c) for c in elem)
        
        try:
            return convert(omath_elem)
        except Exception as e:
            logger.warning(f"OMML conversion failed: {e}")
            # Fallback
            texts = [t.text for t in omath_elem.findall(f'.//{OMML_NS}t') if t.text]
            if texts:
                return f'<span class="formula-fallback">{esc("".join(texts))}</span>'
            return ""

    def _extract_math_segments_from_textframe(self, text_frame) -> list[str]:
        """Извлекает сегменты текста и отдельные формулы (OMML) из текстового фрейма."""
        NS = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "a14": "http://schemas.microsoft.com/office/drawing/2014/main",
            "m": "http://schemas.openxmlformats.org/officeDocument/2006/math",
            "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
        }
        all_paragraphs = []
        for para in text_frame.paragraphs:
            # Обнаружение буллитов
            is_legal_bullet = False
            if para.level > 0:
                is_legal_bullet = True
            else:
                pPr = para._element.pPr
                if pPr is not None:
                    has_bu = any(child.tag.endswith(('buChar', 'buAutoNum', 'buBlip', 'buBullet')) 
                                for child in pPr)
                    if has_bu:
                        is_legal_bullet = True
            
            p_elem = para._element
            segments = []
            current_text = []

            def extract_mathml_from_element(elem):
                """Извлекает MathML из OMML-элемента или возвращает None."""
                mathml = self._omml_to_mathml(elem)
                if mathml:
                    return f"[[[MML_START]]]{mathml}[[[MML_END]]]"
                else:
                    linear_parts = [t.text for t in elem.findall(".//m:t", NS) if t.text]
                    if linear_parts:
                        return f"[[[MML_FB_START]]]{''.join(linear_parts)}[[[MML_FB_END]]]"
                return None

            for child in p_elem:
                local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
                
                # 1. Обработка обычных текстовых runs
                if local == "r":
                    # Проверяем на OMML внутри run
                    omath = child.find(".//m:oMath", NS) or child.find(".//m:oMathPara", NS)
                    if omath is None:
                        omath = child.find(".//mc:AlternateContent/mc:Choice//m:oMath", NS) or \
                                child.find(".//mc:AlternateContent/mc:Choice//m:oMathPara", NS)
                    if omath is not None:
                        if current_text:
                            segments.append("".join(current_text))
                            current_text = []
                        result = extract_mathml_from_element(omath)
                        if result:
                            segments.append(result)
                        continue
                    
                    # Обычный текст + проверка на форматирование (sub/sup/bold/italic)
                    t_elems = child.findall(".//a:t", NS)
                    txt = "".join(t.text or "" for t in t_elems)
                    if txt:
                        rPr = child.find(".//a:rPr", NS)
                        if rPr is not None:
                            bl = rPr.get("baseline")
                            # Поддержка различных способов представления индексов
                            is_sub = (rPr.get("subscript") in ("1", "true", True) or 
                                      child.find(".//a:subscript", NS) is not None)
                            is_sup = (rPr.get("superscript") in ("1", "true", True) or 
                                      child.find(".//a:superscript", NS) is not None)
                            is_bold = rPr.get("b") in ("1", "true", True)
                            is_italic = rPr.get("i") in ("1", "true", True)
                            
                            if bl:
                                try:
                                    val = int(bl)
                                    if val < 0: is_sub = True
                                    elif val > 0: is_sup = True
                                except: pass
                            
                            # Применяем маркеры, которые выживут при экранировании html.escape
                            if is_sub:
                                txt = f"[[SUB_S]]{txt}[[SUB_E]]"
                            elif is_sup:
                                txt = f"[[SUP_S]]{txt}[[SUP_E]]"
                            if is_bold:
                                txt = f"[[B_S]]{txt}[[B_E]]"
                            if is_italic:
                                txt = f"[[I_S]]{txt}[[I_E]]"
                        current_text.append(txt)
                
                # 2. Обработка переноса строки
                elif local == "br":
                    current_text.append("\n")
                
                # 3. Обработка контейнера <a14:m> (ключевое!)
                elif local == "m" and child.tag == f"{{{NS['a14']}}}m":
                    # Ищем OMML внутри <a14:m>
                    omath = child.find(".//m:oMath", NS) or child.find(".//m:oMathPara", NS)
                    if omath is not None:
                        if current_text:
                            segments.append("".join(current_text))
                            current_text = []
                        result = extract_mathml_from_element(omath)
                        if result:
                            segments.append(result)
                    # Пропускаем сам элемент, т.к. содержимое уже обработано
                    continue
                
                # 4. Обработка прямого OMML (на всякий случай)
                elif local in ("oMath", "oMathPara"):
                    if current_text:
                        segments.append("".join(current_text))
                        current_text = []
                    result = extract_mathml_from_element(child)
                    if result:
                        segments.append(result)
                    continue
                
                # 5. Неизвестные элементы – рекурсивно ищем в них OMML
                else:
                    omath = child.find(".//m:oMath", NS) or child.find(".//m:oMathPara", NS)
                    if omath is not None:
                        if current_text:
                            segments.append("".join(current_text))
                            current_text = []
                        result = extract_mathml_from_element(omath)
                        if result:
                            segments.append(result)

            if current_text:
                segments.append("".join(current_text))
            
            # Добавление маркера для буллитов
            if segments:
                if is_legal_bullet:
                    first_seg = segments[0]
                    if not first_seg.startswith("[[[") and not first_seg.startswith("\0"):
                        bullet_chars = set(DESIGN_CONFIG.get("bullet_lists", {}).get("icon_map", {}).keys())
                        stripped = first_seg.lstrip()
                        if not (stripped and stripped[0] in bullet_chars):
                            segments[0] = "• " + first_seg
                all_paragraphs.append(segments)
        
        return all_paragraphs

    def _paragraph_to_html(self, paragraph) -> str:
        """Convert a paragraph (from table cell or shape) to HTML, preserving runs formatting and embedded formulas."""
        NS = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "a14": "http://schemas.microsoft.com/office/drawing/2014/main",
            "m": "http://schemas.openxmlformats.org/officeDocument/2006/math",
            "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
        }
        p_elem = paragraph._element
        parts = []
        for child in p_elem:
            local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if local == "r":
                # Check if this run contains an OMML formula (regular or via mc:AlternateContent)
                omath = child.find(".//m:oMath", NS) or child.find(".//m:oMathPara", NS)
                if omath is None:
                    # Scan deep for AlternateContent
                    omath = child.find(".//mc:AlternateContent/mc:Choice//m:oMath", NS) or \
                            child.find(".//mc:AlternateContent/mc:Choice//m:oMathPara", NS)
                
                if omath is not None:
                    mathml = self._omml_to_mathml(omath)
                    if mathml:
                        # FIXED: Remove redundant <math> wrapper as _omml_to_mathml already provides it
                        parts.append(f'<span class="formula-container">{mathml}</span>')
                    else:
                        # Fallback: linear Unicode text from OMML
                        linear_parts = [t.text for t in omath.findall(".//m:t", NS) if t.text]
                        linear_text = "".join(linear_parts)
                        if linear_text:
                            parts.append(f'<span class="formula-fallback">{html.escape(linear_text)}</span>')
                    continue

                # Regular text run
                t_elems = child.findall(".//a:t", NS)
                txt = "".join(t.text or "" for t in t_elems)
                if not txt:
                    continue
                rPr = child.find(".//a:rPr", NS)
                sub = False
                sup = False
                bold = False
                italic = False
                if rPr is not None:
                    baseline = rPr.get("baseline")
                    if baseline:
                        try:
                            val = int(baseline)
                            if val > 0: sup = True
                            elif val < 0: sub = True
                        except: pass
                    
                    if not sub: sub = rPr.get("subscript") in ("1", "true", True)
                    if not sup: sup = rPr.get("superscript") in ("1", "true", True)
                    bold = rPr.get("b") in ("1", "true", True)
                    italic = rPr.get("i") in ("1", "true", True)
                
                content = html.escape(txt)
                if sub: content = f"<sub>{content}</sub>"
                elif sup: content = f"<sup>{content}</sup>"
                if bold: content = f"<strong>{content}</strong>"
                if italic: content = f"<em>{content}</em>"
                parts.append(content)

            elif local == "br":
                parts.append("<br/>")
            elif local in ("endParaRPr", "extLst"):
                continue
            else:
                omath = child.find(".//m:oMath", NS) or child.find(".//m:oMathPara", NS)
                if omath is not None:
                    mathml = self._omml_to_mathml(omath)
                    if mathml:
                        # FIXED: Remove redundant <math> wrapper
                        parts.append(f'<span class="formula-container">{mathml}</span>')
                    else:
                        linear_parts = [t.text for t in omath.findall(".//m:t", NS) if t.text]
                        linear_text = "".join(linear_parts)
                        if linear_text:
                            parts.append(f'<span class="formula-fallback">{html.escape(linear_text)}</span>')
        return "".join(parts)

    def _table_to_html(self, table) -> str:
        """Преобразует таблицу из PPTX в HTML с сохранением форматирования (case, sub/sup, bold/italic) и формул."""
        rows_html = []
        for row_idx, row in enumerate(table.rows):
            cells_html = ""
            for cell in row.cells:
                # Build cell HTML from all paragraphs
                cell_parts = []
                for para in cell.text_frame.paragraphs:
                    cell_parts.append(self._paragraph_to_html(para))
                cell_content = " ".join(cell_parts)  # separate paragraphs with space
                tag = "th" if row_idx == 0 else "td"
                cells_html += f"<{tag}>{cell_content}</{tag}>"
            rows_html.append(f"<tr>{cells_html}</tr>")
        return f'<table class="data-table"><tbody>{"".join(rows_html)}</tbody></table>'

    def _clean_speaker_name(self, name: str) -> tuple[str, str]:
        """Очищает и разделяет имя докладчика и дополнительную информацию.
        Учитывает наличие маркеров форматирования.
        """
        # Сначала убираем вообще все маркеры для поиска меток типа "Докладчик"
        clean_text = re.sub(r"\[\[.*?\]\]", "", name).strip()
        
        # Удаляем типовые метки
        clean_text = re.sub(r"^[Дд]окладчик:\s*", "", clean_text)
        clean_text = re.sub(r"^[Дд]окладчик\s*[-–—:]\s*", "", clean_text)
        clean_text = re.sub(r"^[Аа]втор:\s*", "", clean_text)
        clean_text = re.sub(r"^[Вв]ыполнил:\s*", "", clean_text)
        
        # Если после очистки ничего не осталось, возвращаем пустые строки
        if not clean_text:
            return "", ""

        parts = clean_text.split("\n")
        if len(parts) >= 2:
            return parts[0].strip(), "\n".join(parts[1:]).strip()
        return clean_text, ""

    def _iter_slide_shapes(self, slide):
        """
        Итерируется по всем фигурам слайда, включая те, что скрыты в mc:AlternateContent.
        Возвращает кортежи (shape_obj, xml_element). 
        Если фигура стандартная, shape_obj — объект BaseShape, xml_element — его _element.
        Если фигура скрыта в AlternateContent, shape_obj может быть None.
        """
        NS = {
            "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
            "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
        }
        
        # Сначала собираем все стандартные фигуры и их ID
        standard_shapes = {}
        for s in slide.shapes:
            try:
                standard_shapes[s.shape_id] = s
            except:
                pass
        
        # Получаем XML дерево слайда
        spTree = slide.shapes._spTree
        
        def process_container(container):
            for child in container:
                # Если это фигура
                tag = child.tag.split('}')[-1]
                if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
                    # Проверяем, есть ли она в стандартном списке
                    # Для этого ищем cNvPr и его id
                    cNvPr = child.find('.//p:cNvPr', {"p": NS["p"]})
                    s_id = int(cNvPr.get('id')) if cNvPr is not None and cNvPr.get('id') else None
                    
                    if s_id in standard_shapes:
                        yield standard_shapes[s_id], child
                    else:
                        yield None, child
                
                # Если это AlternateContent
                elif tag == 'AlternateContent':
                    choice = child.find(f".//{{{NS['mc']}}}Choice")
                    if choice is not None:
                        yield from process_container(choice)
                
                # Если это группа (рекурсивно)
                elif tag == 'grpSp':
                    yield from process_container(child)

        yield from process_container(spTree)

    def extract_content(self) -> None:
        """Извлекает все данные из PowerPoint-файла в self.slides_data."""
        logger.info(f"Чтение презентации: {self.ppt_path}")
        prs = Presentation(self.ppt_path)
        media_output_full = DESIGN_CONFIG["paths"]["media_output_full"]
        os.makedirs(media_output_full, exist_ok=True)

        for i, slide in enumerate(prs.slides):
            slide_info = {
                "title": "",
                "layout_type": "default",
                "content_items": [],  # Unified list for interleaving: {type: 'text'|'table'|'formula', data: ..., pos: (y, x)}
                "visuals": [],
                "is_active": True,
                "slide_num": i + 1,
            }
       
            if slide.shapes.title:
                slide_info["title"] = slide.shapes.title.text

            # Общий проход по всем фигурам (включая скрытые)
            for shape, elem in self._iter_slide_shapes(slide):
                pos = (0, 0)
                if shape:
                    try:
                        pos = (shape.top, shape.left)
                    except: pass
                
                # Пропускаем заголовок слайда (он обрабатывается отдельно)
                if shape == slide.shapes.title:
                    continue

                # 1. СТАНДАРТНЫЙ ТЕКСТ (Text Frame)
                if shape and shape.has_text_frame:
                    try:
                        all_paragraphs = self._extract_math_segments_from_textframe(shape.text_frame)
                        if all_paragraphs:
                            slide_info["content_items"].append({
                                "type": "text",
                                "data": all_paragraphs,
                                "pos": pos
                            })
                            # Статистика
                            for para in all_paragraphs:
                                f_count = sum(1 for s in para if isinstance(s, str) and (s.startswith("[[[MML_START]]]") or s.startswith("[[[MML_FB_START]]]")))
                                self.stats["formulas"] += f_count
                    except Exception as e:
                        logger.warning(f"Slide {i + 1}: error extracting text from shape: {e}")

                # 2. ТАБЛИЦЫ
                elif shape and shape.has_table:
                    try:
                        table_html = self._table_to_html(shape.table)
                        slide_info["content_items"].append({
                                "type": "table",
                                "data": table_html,
                                "pos": pos
                        })
                        self.stats["tables"] += 1
                    except Exception as e:
                        logger.error(f"Slide {i + 1}: table error: {e}")

                # 3. ФОРМУЛЫ (Standalone OMML / OLE с OMML)
                else:
                    NS_M = {"m": "http://schemas.openxmlformats.org/officeDocument/2006/math"}
                    # Explicit check to avoid FutureWarning in truth value of elements
                    omath = elem.find('.//m:oMath', NS_M)
                    if omath is None:
                        omath = elem.find('.//m:oMathPara', NS_M)
                    
                    if omath is not None:
                        try:
                            # Пытаемся извлечь позицию из XML, если shape=None
                            if not shape:
                                off = elem.find('.//a:off', {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
                                if off is not None:
                                    pos = (int(off.get('y', 0)), int(off.get('x', 0)))
                            
                            mathml = self._omml_to_mathml(omath)
                            if mathml:
                                slide_info["content_items"].append({
                                    "type": "formula",
                                    "data": mathml,
                                    "pos": pos
                                })
                                self.stats["formulas"] += 1
                        except Exception as e:
                            logger.warning(f"Slide {i + 1}: standalone math extraction failed: {e}")

                # 4. ИЗОБРАЖЕНИЯ (Остаются в visuals для отдельного рендеринга справа)
                if shape and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_name = f"slide_{i + 1}_img_{len(slide_info['visuals']) + 1}.png"
                    img_path_full = str(Path(media_output_full) / img_name)
                    img_path_rel = (
                        f"{DESIGN_CONFIG['paths']['media_output']}/{img_name}"
                    )

                    try:
                        try:
                            blob = shape.image.blob
                        except ValueError:
                            if not HAS_LXML:
                                logger.warning(
                                    "lxml не установлен, нельзя извлечь SVG. Пропускаем изображение."
                                )
                                self.stats["images_fail"] += 1
                                continue

                            elem = shape._element
                            svg_blips = elem.findall(
                                ".//{http://schemas.microsoft.com/office/drawing/2016/SVG/main}svgBlip"
                            )
                            if not svg_blips:
                                logger.warning(
                                    f"SVG blip не найден для slide {i + 1}, пропускаем"
                                )
                                self.stats["images_fail"] += 1
                                continue
                            rId = svg_blips[0].get(
                                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                            )
                            if not rId:
                                logger.warning(
                                    f"Нет rId для SVG blip на слайде {i + 1}"
                                )
                                self.stats["images_fail"] += 1
                                continue
                            rel = slide.part.rels[rId]
                            part = rel.target_part
                            svg_blob = part.blob
                            if not svg_blob or len(svg_blob) == 0:
                                logger.warning(
                                    f"Пустой SVG blob на слайде {i + 1}, пропускаем"
                                )
                                if os.path.exists(img_path_full):
                                    try:
                                        os.remove(img_path_full)
                                    except Exception:
                                        pass
                                self.stats["images_fail"] += 1
                                continue
                            svg_path_full = img_path_full.replace(".png", ".svg")
                            svg_path_rel = img_path_rel.replace(".png", ".svg")

                            try:
                                import cairosvg

                                cairosvg.svg2png(
                                    bytestring=svg_blob,
                                    write_to=img_path_full,
                                    background_color="white",
                                )
                                if (
                                    os.path.exists(img_path_full)
                                    and os.path.getsize(img_path_full) > 0
                                ):
                                    slide_info["visuals"].append(
                                        {
                                            "src": img_path_rel,
                                            "caption": "",
                                            "pos": (
                                                shape.left,
                                                shape.top,
                                                shape.width,
                                                shape.height,
                                            ),
                                        }
                                    )
                                    self.stats["images_ok"] += 1
                                    continue
                                else:
                                    raise IOError("Converted PNG is empty")
                            except Exception as svg_err:
                                logger.warning(
                                    f"cairosvg не сработал ({svg_err}), сохраняем SVG как есть"
                                )
                                if os.path.exists(img_path_full):
                                    try:
                                        os.remove(img_path_full)
                                    except Exception:
                                        pass
                                with open(svg_path_full, "wb") as sf:
                                    sf.write(svg_blob)
                                slide_info["visuals"].append(
                                    {
                                        "src": svg_path_rel,
                                        "caption": "",
                                        "pos": (
                                            shape.left,
                                            shape.top,
                                            shape.width,
                                            shape.height,
                                        ),
                                    }
                                )
                                self.stats["images_ok"] += 1
                                continue

                        if not blob or len(blob) == 0:
                            logger.warning(
                                f"Пустой blob изображения на слайде {i + 1}, пропускаем"
                            )
                            if os.path.exists(img_path_full):
                                try:
                                    os.remove(img_path_full)
                                except Exception:
                                    pass
                            self.stats["images_fail"] += 1
                            continue

                        success = self._save_image_with_white_bg(blob, img_path_full)
                        if success:
                            slide_info["visuals"].append(
                                {
                                    "src": img_path_rel,
                                    "caption": "",
                                    "pos": (
                                        shape.left,
                                        shape.top,
                                        shape.width,
                                        shape.height,
                                    ),
                                }
                            )
                            self.stats["images_ok"] += 1
                        else:
                            self.stats["images_fail"] += 1
                    except Exception as e:
                        logger.error(f"[IMG] Ошибка на слайде {i + 1}: {e}")
                        self.stats["images_fail"] += 1

                # OLE objects support
                else:
                    try:
                        if shape and hasattr(shape, "ole_format") and shape.ole_format is not None:
                            # Deep scanning for OMML already covered by search loop
                            pass
                    except (ValueError, AttributeError):
                        pass

            # Применяем строгую пространственную сортировку ко всему контенту
            slide_info["content_items"] = self._spatial_sort_strict(slide_info["content_items"])
            # Сортировка изображений (остается стандартной)
            slide_info["visuals"] = self._spatial_sort(slide_info["visuals"])

            # Очистка и доп. логика для интро-слайда
            if i == 0:
                slide_info["layout_type"] = "intro"
                presenter_items = []
                for item in slide_info["content_items"]:
                    if item["type"] == "text":
                        for p in item["data"]:
                            txt = "".join(s for s in p if not s.startswith("[[[") ).strip()
                            if txt: presenter_items.append(txt)
                
                if presenter_items:
                    # Попытаемся найти именно имя, пропуская явные заголовки "Докладчик"
                    title_candidate = presenter_items[0]
                    speaker_candidate = ""
                    info_candidate = ""
                    
                    # Просматриваем остальные элементы
                    remaining = presenter_items[1:]
                    found_name = False
                    
                    for i, item in enumerate(remaining):
                        pure_item = re.sub(r"\[\[.*?\]\]", "", item).strip()
                        # Если это просто слово "Докладчик", пропускаем его и берем следующее как имя
                        if pure_item.lower() in ("докладчик", "докладчик:", "выполнил", "автор"):
                            continue
                        
                        if not found_name:
                            name, info = self._clean_speaker_name("\n".join(remaining[i:]))
                            speaker_candidate = name
                            info_candidate = info
                            found_name = True
                            break
                    
                    # Фолбэк, если ничего не нашли по логике выше
                    if not found_name and remaining:
                        name, info = self._clean_speaker_name("\n".join(remaining))
                        speaker_candidate = name
                        info_candidate = info

                    slide_info["title"] = title_candidate
                    slide_info["speaker_name"] = speaker_candidate if speaker_candidate else "Докладчик"
                    slide_info["speaker_info"] = info_candidate
                else:
                    slide_info["speaker_name"] = "Докладчик"
                    slide_info["speaker_info"] = ""

            elif len(slide_info["visuals"]) == 2:
                slide_info["layout_type"] = "two_images"
            elif len(slide_info["visuals"]) == 0:
                # Проверяем наличие таблиц в content_items
                has_table = any(it["type"] == "table" for it in slide_info["content_items"])
                if not has_table:
                    slide_info["layout_type"] = "full_text"

            if slide_info["visuals"]:
                self._extract_captions_from_shapes(slide, slide_info)

            # Логирование формул
            f_total = sum(1 for it in slide_info["content_items"] if it["type"] == "formula")
            # Считаем встроенные формулы
            for it in slide_info["content_items"]:
                if it["type"] == "text":
                    for p in it["data"]:
                        f_total += sum(1 for s in p if isinstance(s, str) and s.startswith("[[[MML"))
            
            if f_total > 0:
                logger.info(f"Slide {slide_info['slide_num']}: total {f_total} formulas found.")

            self.slides_data.append(slide_info)

        self.stats["total_slides"] = len(self.slides_data)

    def _generate_section_tag(self, data: dict) -> str:
        """Генерирует тег секции (например, "Данные", "Результаты") на основе текста."""
        all_texts = []
        for block in data.get("content_items", []):
            if block["type"] == "text":
                for p in block["data"]:
                    all_texts.append(" ".join(s for s in p if not s.startswith("[[[") ))
        
        text_combined = " ".join(all_texts).lower()
        if any(kw in text_combined for kw in ["данные", "таблиц", "исходн"]):
            return "Данные"
        if any(kw in text_combined for kw in ["результат", "вывод", "итог"]):
            return "Результаты"
        if any(kw in text_combined for kw in ["описание", "метод", "подход"]):
            return "Описание"
        if any(kw in text_combined for kw in ["анализ", "исследовани"]):
            return "Анализ"
        return "Описание"

    def _format_text_panel(self, slide_info: dict) -> str:
        """Формирует HTML-панель с чередованием текста, таблиц и формул."""
        parts = []
        bullet_config = DESIGN_CONFIG.get("bullet_lists", {})
        bullet_icon_map = bullet_config.get("icon_map", {})
        bullet_icon_size = DESIGN_CONFIG["icon_size"]
        bullet_indent = bullet_config.get("indent", "2rem")
        bullet_border = bullet_config.get("border_left", "2px solid rgba(0,242,255,0.15)")
        bullet_bg = bullet_config.get("background", "rgba(0,242,255,0.03)")

        # Итерируемся по объединенному списку контента для соблюдения порядка
        for item_block in slide_info.get("content_items", []):
            item_type = item_block["type"]
            data = item_block["data"]

            if item_type == "text":
                # Обработка сегментированных параграфов
                for p in data:
                    # 1. Сбор сегментов и замена формул на токены
                    segments_to_join = []
                    formula_store = {}
                    
                    if isinstance(p, list):
                        for idx, seg in enumerate(p):
                            if seg.startswith("[[[MML_START]]]"):
                                token = f"[[_F{idx}_]]"
                                mathml = seg[len("[[[MML_START]]]"): -len("[[[MML_END]]]")]
                                formula_store[token] = f'<span class="formula-container">{mathml}</span>'
                                segments_to_join.append(token)
                            elif seg.startswith("[[[MML_FB_START]]]"):
                                token = f"[[_FB{idx}_]]"
                                fb = seg[len("[[[MML_FB_START]]]"): -len("[[[MML_FB_END]]]")]
                                formula_store[token] = f'<span class="formula-fallback">{esc(fb)}</span>'
                                segments_to_join.append(token)
                            else:
                                segments_to_join.append(seg)
                    else:
                        segments_to_join.append(str(p))

                    full_content_with_tokens = "".join(segments_to_join)
                    if not full_content_with_tokens.strip():
                        continue

                    # 2. Очистка для поиска буллитов (наши токены выживут)
                    clean_search_text = re.sub(r"<[^>]+>", "", full_content_with_tokens)
                    items = self._split_text_into_items(clean_search_text)

                    for item in items:
                        # 3. Экранируем текст и возвращаем формулы и форматирование
                        display_text = html.escape(item["text"])
                        
                        # Возврат форматирования из маркеров
                        display_text = display_text.replace("[[SUB_S]]", "<sub>").replace("[[SUB_E]]", "</sub>")
                        display_text = display_text.replace("[[SUP_S]]", "<sup>").replace("[[SUP_E]]", "</sup>")
                        display_text = display_text.replace("[[B_S]]", "<strong>").replace("[[B_E]]", "</strong>")
                        display_text = display_text.replace("[[I_S]]", "<em>").replace("[[I_E]]", "</em>")

                        for token, real_html in formula_store.items():
                            display_text = display_text.replace(token, real_html)
                        
                        if item.get("is_bullet"):
                            keyword_icon = self.get_icon_by_text(item["text"])
                            default_bullet = "chevron-right"
                            if keyword_icon != default_bullet:
                                icon = keyword_icon
                            else:
                                icon = bullet_icon_map.get(item.get("bullet_char"), "chevron-right")

                            parts.append(
                                f'<div class="list-item-bullet" style="padding-left: {bullet_indent}; border-left: {bullet_border}; background: {bullet_bg};">'
                                f'<i data-lucide="{icon}" style="width: {bullet_icon_size}; height: {bullet_icon_size}; flex-shrink: 0;"></i>'
                                f'<div class="list-text">{display_text}</div></div>'
                            )
                        else:
                            keyword_icon = self.get_icon_by_text(item["text"])
                            parts.append(
                                f'<div class="list-item">'
                                f'<i data-lucide="{keyword_icon}" style="width: {DESIGN_CONFIG["icon_size"]}; height: {DESIGN_CONFIG["icon_size"]}; flex-shrink: 0;"></i>'
                                f'<div class="list-text">{display_text}</div></div>'
                            )

            elif item_type == "table":
                parts.append(str(data))

            elif item_type == "formula":
                parts.append(f'<div class="formula-block animate-up">{data}</div>')

        return "".join(parts)
        

    def generate_html(self) -> None:
        """Генерирует итоговый HTML-файл на основе встроенного шаблона и данных слайдов."""
        logger.info("Рендеринг HTML...")

        # Регулировка <title> и имени файла на основе данных докладчика
        speaker_name = self.slides_data[0].get("speaker_name", "") if self.slides_data else ""
        # Очистка имени от ВСЕХ маркеров форматирования для использования в заголовках и именах файлов
        clean_speaker = re.sub(r"\[\[.*?\]\]", "", speaker_name)
        
        # 1. Сначала определяем имя на основе докладчика
        if clean_speaker:
            safe_name = re.sub(r'[\\/*?:"<>|]', '', clean_speaker).strip()
            if safe_name:
                # Добавляем ПНИПУ только в название файла
                self.output_html = f"{safe_name} (ПНИПУ).html"
        
        # 2. Если имя все еще не определено (и не было передано снаружи), используем внутренний фолбэк
        if not self.output_html:
            self.output_html = "presentation_output.html"

        # --- Standalone / Embedding Mode ---
        is_standalone = DESIGN_CONFIG.get("STATIC_ASSETS_EMBED", False)
        head_part = BASE_HTML_TEMPLATE
        tail_part = BASE_HTML_TAIL
        
        # Регулировка <title> в head_part
        s_title = clean_speaker if clean_speaker else "Доклад"
        head_part = head_part.replace("{speaker_name}", s_title)

        # Logo path resolution
        logo_rel = DESIGN_CONFIG["paths"]["logo_white"]
        v_logo_path = BASE_DIR / "web_demo" / logo_rel
        if not v_logo_path.exists():
            v_logo_path = BASE_DIR / logo_rel
            
        if is_standalone:
            logger.info("Подготовка автономного (standalone) файла (все ресурсы вшиваются)...")
            # 1. Инлайним логотип
            logo = self._get_data_uri(v_logo_path)
            
            # 2. Инлайним шрифты
            fonts_css_path = BASE_DIR / "libs" / "fonts" / "fonts.css"
            fonts_inlined = self._inline_css_fonts(fonts_css_path)
            head_part = re.sub(
                r'<link rel="stylesheet" href="libs/fonts/fonts.css">',
                f'<style>{fonts_inlined}</style>',
                head_part
            )
            
            # 3. Инлайним скрипты
            scripts_to_inline = [
                ("libs/gsap/gsap.min.js", r'<script src="libs/gsap/gsap.min.js" defer></script>'),
                ("libs/lucide/lucide.min.js", r'<script src="libs/lucide/lucide.min.js" defer></script>'),
                ("libs/mathjax/tex-mml-svg.js", r'<script src="libs/mathjax/tex-mml-svg.js" defer></script>')
            ]
            
            for rel_path, pattern in scripts_to_inline:
                s_path = BASE_DIR / rel_path
                s_content = self._get_file_content(s_path)
                # Robust replacement using lambda to avoid escape issues in large blocks
                head_part = re.sub(pattern, lambda m, c=s_content: f'<script>{c}</script>', head_part, flags=re.DOTALL)
        else:
            logo = DESIGN_CONFIG["paths"]["logo_white"]

        slides_content = ""
        total = len(self.slides_data)

        for idx, data in enumerate(self.slides_data):
            num = idx + 1
            title = data["title"]
            section_tag = self._generate_section_tag(data)
            text_panel_html = self._format_text_panel(data)

            # --- Умный шрифт: расчет panel_class на основе сложности контента ---
            total_text_chars = 0
            for item in data.get("content_items", []):
                if item["type"] == "text":
                    for p in item["data"]:
                        for s in p:
                            if isinstance(s, str):
                                # Считаем количество формул в сегменте
                                f_count = s.count("[[[MML_START]]]") + s.count("[[[MML_FB_START]]]")
                                if f_count > 0:
                                    # Очищаем сегмент от кода формул для честного подсчета текста
                                    # Ищем все вхождения [[[MML...]]] и удаляем их содержимое
                                    clean_s = re.sub(r"\[\[\[MML_START\]\].*?\[\[\[MML_END\]\]\]", "", s, flags=re.DOTALL)
                                    clean_s = re.sub(r"\[\[\[MML_FB_START\]\].*?\[\[\[MML_FB_END\]\]\]", "", clean_s, flags=re.DOTALL)
                                    # Также очищаем от новых маркеров форматирования
                                    clean_s = re.sub(r"\[\[[A-Z0-9_]+\]\]", "", clean_s)
                                    
                                    total_text_chars += len(clean_s)
                                    total_text_chars += (f_count * 20) # Significantly reduced weight for inline formulas
                                else:
                                    clean_s = re.sub(r"\[\[[A-Z0-9_]+\]\]", "", s)
                                    total_text_chars += len(clean_s)
                elif item["type"] == "table":
                    total_text_chars += 300  # Reduced table density impact
                elif item["type"] == "formula":
                    total_text_chars += 40 # Standalone formulas weight even less now


            panel_class = "analytical-panel animate-up"
            # -------------------------------------------------------------

            if is_standalone:
                embedded_count = 0
                for vis in data.get("visuals", []):
                    if vis.get("src") and not vis["src"].startswith("data:"):
                        # Картинки в web_demo/media/... относительны корня проекта в slides_data
                        v_path = BASE_DIR / "web_demo" / vis["src"]
                        if v_path.exists():
                            vis["src"] = self._get_data_uri(v_path)
                            embedded_count += 1
                        else:
                            logger.warning(f"Медиа не найдено для эмбеддинга: {v_path}")
                if embedded_count > 0:
                    logger.debug(f"Слайд {idx+1}: эмбедировано {embedded_count} изображений")
            
            if data["layout_type"] == "intro":
                slides_content += f"""
        <section class="slide hide-title">
            <div class="logo-container"><img src="{logo}" alt="Логотип" class="header-logo" onerror="this.style.display='none'; this.onerror=null;"></div>
            <div class="slide-header">
                <div class="slide-number">{num:02d} / {total:02d}</div>
            </div>
            <div class="slide-content-title">
                <h1 class="main-heading animate-up">{esc(title)}</h1>
            </div>
            <div class="presenter-card animate-up">
                <span class="presenter-label">Докладчик</span>
                <div class="presenter-name">{esc(data.get("speaker_name", ""))}</div>
                <div class="presenter-info">{esc(data.get("speaker_info", ""))}</div>
            </div>
        </section>"""

            elif data["layout_type"] == "conclusions_dual":
                slides_content += f"""
        <section class="slide">
            <div class="logo-container"><img src="{logo}" alt="Логотип" class="header-logo" onerror="this.style.display='none'; this.onerror=null;"></div>
            <div class="slide-header"><div class="slide-number">{num:02d} / {total:02d}</div><div class="slide-title">{esc(title)}</div></div>
            <div class="slide-split" style="grid-template-columns: 1fr 1fr; height: calc(100% - 145px);">
                <div class="{panel_class}">
                    {data["left_html"]}
                </div>
                <div class="{panel_class}">
                    {data["right_html"]}
                </div>
            </div>
        </section>"""

            elif data["layout_type"] == "research_roadmap":
                slides_content += f"""
        <section class="slide">
            <div class="logo-container"><img src="{logo}" alt="Логотип" class="header-logo" onerror="this.style.display='none'; this.onerror=null;"></div>
            <div class="slide-header"><div class="slide-number">{num:02d} / {total:02d}</div><div class="slide-title">{esc(title)}</div></div>
            <div class="slide-split" style="grid-template-columns: 1.4fr 0.8fr; height: calc(100% - 145px);">
                <div class="{panel_class}">
                    <span class="section-tag" style="font-size: var(--fs-tag);">План 2026</span>
                    {data["content_html"]}
                </div>
                <div class="viz-card animate-up" style="display: flex; flex-direction: column; justify-content: center; align-items: center; background: radial-gradient(circle, var(--accent-soft) 0%, transparent 80%); border-radius: 3rem; padding: 3rem; border: 1px solid var(--glass-border); position: relative; overflow: hidden; height: 100%;">
                     <div class="rocket-glow"></div>
                     <i data-lucide="rocket" style="width: 120px; height: 120px; color: var(--accent); margin-bottom: 2rem; filter: drop-shadow(0 0 30px var(--accent)); transform: rotate(-45deg); animation: pulse-rocket 2s infinite ease-in-out;"></i>
                     <h4 style="font-family: 'Outfit'; font-size: var(--fs-research-year); margin: 0; font-weight: 800; background: linear-gradient(135deg, white 0%, var(--accent) 100%); -webkit-background-clip: text; -webkit-text-fill-color: transparent;">ROADMAP</h4>
                </div>
            </div>
        </section>"""

            elif data["layout_type"] == "full_text":
                slides_content += f"""
        <section class="slide">
            <div class="logo-container"><img src="{logo}" alt="Логотип" class="header-logo" onerror="this.style.display='none'; this.onerror=null;"></div>
            <div class="slide-header"><div class="slide-number">{num:02d} / {total:02d}</div><div class="slide-title">{esc(title)}</div></div>
            <div class="slide-content-full animate-up" style="width: 85%; margin: 2rem auto; height: calc(100% - 145px); min-height: 0;">
                <div class="{panel_class}" style="height: 100%;">
                    {text_panel_html}
                </div>
            </div>
        </section>"""

            else:
                # Стандартный макет со сплитом или только изображениями
                
                # Сортировка визуальных элементов
                data["visuals"] = self._spatial_sort(data["visuals"])
                
                # --- Расчет динамических пропорций (Dynamic Ratio) ---
                has_text = bool(text_panel_html.strip())
                base_ratio = DESIGN_CONFIG["layout"]["text_panel_ratio"] # 0.35
                layout_class = "slide-split"
                
                # 1. Предварительный расчет компоновки для определения количества колонок
                # Мы делаем "пробный" запуск с условной шириной, чтобы понять, пойдут ли картинки в один ряд или в несколько
                pre_results = self.get_best_layout(data["visuals"], 1200, 940)
                rows_pre, cols_pre = pre_results[0], pre_results[1]
                
                if has_text:
                    if cols_pre == 1:
                        # --- Адаптивный режим для ОДНОЙ КОЛОНКИ рисунков (1 шт или стек по вертикали) ---
                        grid_split = "1fr auto"
                        layout_class += " layout-auto-width"
                        cont_w = 1200 # Условное значение для расчета Area
                    else:
                        # --- Расчет по плотности для МНОГОКОЛОНОЧНЫХ компоновок ---
                        d_factor = max(0.65, min(1.2, total_text_chars / 1200.0 + 0.5))
                        dynamic_ratio = base_ratio * d_factor
                        dynamic_ratio = max(0.25, min(0.42, dynamic_ratio))
                        
                        grid_split = f"{dynamic_ratio:.3f}fr {1.0 - dynamic_ratio:.3f}fr"
                        cont_w = 1728 * (1.0 - dynamic_ratio)
                else:
                    cont_w = 1728
                    grid_split = "1fr"

                cont_h = 940 
                # Финальный расчет с уточненной шириной
                results = self.get_best_layout(data["visuals"], cont_w, cont_h)
                rows, cols, grid_styles, row_tmpl, col_tmpl = results

                # Сборка HTML визуальных элементов
                visuals_items_html = ""
                for idx, vis in enumerate(data["visuals"]):
                    style = grid_styles[idx] if idx < len(grid_styles) else ""
                    caption_html = (
                        f'<div class="viz-caption">{esc(vis.get("caption", ""))}</div>'
                        if vis.get("caption")
                        else ""
                    )
                    if vis.get("src"):
                        alt_text = (esc(vis["caption"]) if vis.get("caption") else "Изображение")
                        content_html = f'<div class="viz-box"><img src="{vis["src"]}" alt="{alt_text}"></div>'
                    else:
                        content_html = '<div class="error-box">Нет изображения</div>'
                    visuals_items_html += f'<div class="viz-item" style="{style}">{caption_html}{content_html}</div>'

                if has_text:
                    slides_content += f"""
        <section class="slide">
            <div class="logo-container"><img src="{logo}" alt="Логотип" class="header-logo" onerror="this.style.display='none'; this.onerror=null;"></div>
            <div class="slide-header"><div class="slide-number">{num:02d} / {total:02d}</div><div class="slide-title">{esc(title)}</div></div>
            <div class="slide-split {layout_class}" style="grid-template-columns: {grid_split}; height: calc(100% - 145px); min-height: 0;">
                <div class="{panel_class}"><span class="section-tag" style="font-size: var(--fs-tag);">{esc(section_tag)}</span>{text_panel_html}</div>
                <div class="img-stack animate-up" style="display: grid; grid-template-columns: {col_tmpl}; grid-template-rows: {row_tmpl}; gap: var(--gap-main);">{visuals_items_html}</div>
            </div>
        </section>"""
                else:
                    # Images only: full-width grid
                    slides_content += f"""
        <section class="slide">
            <div class="logo-container"><img src="{logo}" alt="Логотип" class="header-logo" onerror="this.style.display='none'; this.onerror=null;"></div>
            <div class="slide-header"><div class="slide-number">{num:02d} / {total:02d}</div><div class="slide-title">{esc(title)}</div></div>
            <div class="slide-content-full animate-up" style="width: 100%; margin: 0; height: calc(100% - 145px); display: grid; grid-template-columns: {col_tmpl}; grid-template-rows: {row_tmpl}; gap: var(--gap-main); padding: 0;">
                {visuals_items_html}
            </div>
        </section>"""

        logger.info(f"Сохранение в {self.output_html}...")
        with open(self.output_html, "w", encoding="utf-8") as f:
            f.write(head_part)
            f.write(slides_content)
            f.write(tail_part)

        logger.info("\n" + "=" * 40)
        logger.info(" ИТОГОВЫЙ ОТЧЕТ КОНВЕРТАЦИИ")
        logger.info("=" * 40)
        logger.info(f" Всего слайдов:   {self.stats['total_slides']}")
        logger.info(f" Успешных фото:   {self.stats['images_ok']}")
        logger.info(f" Таблиц:          {self.stats['tables']}")
        logger.info(f" Формул:          {self.stats['formulas']}")
        logger.info(f" Пропущено:       {self.stats['images_fail']} (см. лог выше)")
        logger.info("=" * 40)


if __name__ == "__main__":
    ppt_file = "Промежуточная.pptx"

    if os.path.exists(ppt_file):
        conv = PPTConverter(ppt_file)
        conv.extract_content()
        conv.process_txt_files()
        conv.generate_html()
        logger.info(f"Готово! Файл создан: {conv.output_html}")
    else:
        logger.error(f"Не найден входной файл .pptx: {ppt_file}")

