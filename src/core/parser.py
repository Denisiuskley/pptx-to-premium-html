import requests
import urllib.request
import urllib.parse
import json
import hashlib
from ..utils.llm_cache import LLMCache
import os
import re
import html
import logging
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
from typing import List, Dict, Any, Optional
from ..config import DESIGN_CONFIG, OMML_NS, MATHML_NS
from ..utils.text_helpers import clean_text, _is_slide_number, esc
from ..utils.image_helpers import _save_image_with_white_bg
from ..converters.math_converter import omml_to_mathml
logger = logging.getLogger(__name__)

class PPTParser:
    """Извлекает структурированный контент из файлов PowerPoint, сохраняя иерархию и формулы."""

    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.presentation = Presentation(pptx_path)
        self.slide_width = self.presentation.slide_width
        self.slide_height = self.presentation.slide_height
        self.stats = {'total_slides': 0, 'images_ok': 0, 'images_fail': 0, 'tables': 0, 'formulas': 0}
        self.slides_data = []

    def extract_content(self) -> List[Dict[str, Any]]:
        """Главный цикл обхода слайдов и извлечения данных."""
        logger.info(f'Начало извлечения из: {self.pptx_path}')
        media_output_full = DESIGN_CONFIG['paths']['media_output_full']
        os.makedirs(media_output_full, exist_ok=True)
        for i, slide in enumerate(self.presentation.slides):
            slide_info = {'slide_num': i + 1, 'title': '', 'layout_type': 'default', 'content_items': [], 'visuals': [], 'speaker_name': '', 'speaker_info': ''}
            if slide.shapes.title:
                slide_info['title'] = slide.shapes.title.text
            for shape, elem in self._iter_slide_shapes(slide):
                pos = (0, 0)
                if shape:
                    try:
                        pos = (shape.top, shape.left)
                    except:
                        pass
                if shape == slide.shapes.title:
                    continue
                if shape and shape.has_text_frame:
                    try:
                        shape_txt = clean_text(shape.text)
                        if shape_txt.lower() == slide_info['title'].lower():
                            continue
                        paragraphs = self._extract_math_segments_from_textframe(shape.text_frame)
                        if paragraphs:
                            slide_info['content_items'].append({'type': 'text', 'data': paragraphs, 'pos': pos})
                            for para in paragraphs:
                                f_count = sum((1 for s in para if isinstance(s, str) and ('[[[MML_START]]]' in s or '[[[MML_FB_START]]]' in s)))
                                self.stats['formulas'] += f_count
                    except Exception as e:
                        logger.warning(f'Ошибка парсинга текста на слайде {i + 1}: {e}')
                elif shape and shape.has_table:
                    try:
                        table_html = self._table_to_html(shape.table)
                        slide_info['content_items'].append({'type': 'table', 'data': table_html, 'pos': pos})
                        self.stats['tables'] += 1
                    except Exception as e:
                        logger.error(f'Ошибка таблицы на слайде {i + 1}: {e}')
                if shape and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        self._process_image(shape, i, slide_info)
                    except Exception as e:
                        logger.error(f'Ошибка изображения на слайде {i + 1}: {e}')
                
                # 3. ФОРМУЛЫ (Standalone OMML / OLE с OMML)
                # Ищем OMML в XML-элементе, даже если shape=None или это PICTURE/GraphicFrame
                NS_M = {"m": "http://schemas.openxmlformats.org/officeDocument/2006/math"}
                omath = elem.find('.//m:oMath', NS_M)
                if omath is None:
                    omath = elem.find('.//m:oMathPara', NS_M)
                
                if omath is not None:
                    try:
                        # Если это не текст и не таблица, добавляем как отдельный блок формулы
                        # (В тексте формулы извлекаются через _extract_math_segments_from_textframe)
                        if not (shape and (shape.has_text_frame or shape.has_table)):
                            # Пытаемся извлечь позицию из XML, если shape=None
                            if not shape:
                                off = elem.find('.//a:off', {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
                                if off is not None:
                                    pos = (int(off.get('y', 0)), int(off.get('x', 0)))
                            
                            mathml = self._omml_to_mathml(omath)
                            if mathml:
                                slide_info['content_items'].append({
                                    'type': 'formula',
                                    'data': mathml,
                                    'pos': pos
                                })
                                self.stats['formulas'] += 1
                    except Exception as e:
                        logger.warning(f'Slide {i + 1}: standalone math extraction failed: {e}')
            if slide_info['visuals']:
                self._extract_captions_from_shapes(slide, slide_info)
            slide_info['content_items'].sort(key=lambda x: x['pos'][0])
            slide_info['visuals'].sort(key=lambda x: (x['pos'][1], x['pos'][0]))
            if i == 0:
                slide_info['layout_type'] = 'intro'
                self._enrich_intro_slide(slide_info)
            elif len(slide_info['visuals']) == 0:
                has_table = any((it['type'] == 'table' for it in slide_info['content_items']))
                if not has_table:
                    slide_info['layout_type'] = 'full_text'
            self.slides_data.append(slide_info)
        self.stats['total_slides'] = len(self.slides_data)
        return self.slides_data

    def _iter_slide_shapes(self, slide):
        """Итерируется по всем фигурам слайда, включая скрытые в AlternateContent."""
        NS = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main', 'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006'}
        standard_shapes = {s.shape_id: s for s in slide.shapes if hasattr(s, 'shape_id')}
        spTree = slide.shapes._spTree

        def process_container(container):
            for child in container:
                tag = child.tag.split('}')[-1]
                if tag in ('sp', 'pic', 'graphicFrame', 'grpSp', 'cxnSp'):
                    cNvPr = child.find('.//p:cNvPr', {'p': NS['p']})
                    s_id = int(cNvPr.get('id')) if cNvPr is not None and cNvPr.get('id') else None
                    if s_id in standard_shapes:
                        yield (standard_shapes[s_id], child)
                    else:
                        yield (None, child)
                elif tag == 'AlternateContent':
                    choice = child.find(f".//{{{NS['mc']}}}Choice")
                    if choice is not None:
                        yield from process_container(choice)
                elif tag == 'grpSp':
                    yield from process_container(child)
        yield from process_container(spTree)

    def _extract_math_segments_from_textframe(self, text_frame) -> list:
        """Извлекает сегменты текста и MathML маркеры из текстового фрейма."""
        NS = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main', 'm': 'http://schemas.openxmlformats.org/officeDocument/2006/math', 'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006', 'a14': 'http://schemas.microsoft.com/office/drawing/2014/main'}
        all_paragraphs = []
        for para in text_frame.paragraphs:
            is_bullet = para.level > 0 or (para._element.pPr is not None and any((c.tag.endswith(('buChar', 'buAutoNum', 'buBlip', 'buBullet')) for c in para._element.pPr)))
            p_elem = para._element
            segments = []
            current_text = []
            for child in p_elem:
                local = child.tag.split('}')[-1]
                if local == 'r':
                    omath = child.find('.//m:oMath', NS) or child.find('.//mc:AlternateContent/mc:Choice//m:oMath', NS)
                    if omath is not None:
                        if current_text:
                            segments.append(''.join(current_text))
                            current_text = []
                        segments.append(self._get_mml_token(omath))
                        continue
                    t_elems = child.findall('.//a:t', NS)
                    txt = ''.join((t.text or '' for t in t_elems))
                    if txt:
                        rPr = child.find('.//a:rPr', NS)
                        if rPr is not None:
                            is_sub = rPr.get('baseline') and int(rPr.get('baseline')) < 0 or rPr.get('subscript') in ('1', 'true')
                            is_sup = rPr.get('baseline') and int(rPr.get('baseline')) > 0 or rPr.get('superscript') in ('1', 'true')
                            if is_sub:
                                txt = f'[[SUB_S]]{txt}[[SUB_E]]'
                            elif is_sup:
                                txt = f'[[SUP_S]]{txt}[[SUP_E]]'
                            if rPr.get('b') in ('1', 'true'):
                                txt = f'[[B_S]]{txt}[[B_E]]'
                            if rPr.get('i') in ('1', 'true'):
                                txt = f'[[I_S]]{txt}[[I_E]]'
                        current_text.append(txt)
                elif local == 'br':
                    current_text.append('\n')
                elif local in ('oMath', 'oMathPara'):
                    if current_text:
                        segments.append(''.join(current_text))
                        current_text = []
                    segments.append(self._get_mml_token(child))
            if current_text:
                segments.append(''.join(current_text))
            if segments:
                if is_bullet and segments[0] and (not segments[0].startswith('[[[')):
                    segments[0] = '• ' + segments[0]
                all_paragraphs.append(segments)
        return all_paragraphs

    def _get_mml_token(self, omath_elem) -> str:
        mathml = omml_to_mathml(omath_elem)
        if mathml:
            return f'[[[MML_START]]]{mathml}[[[MML_END]]]'
        texts = [t.text for t in omath_elem.findall(f'.//{OMML_NS}t') if t.text]
        return f"[[[MML_FB_START]]]{''.join(texts)}[[[MML_FB_END]]]" if texts else ''

    def _omml_to_mathml(self, omath_elem) -> Optional[str]:
        """Обертка над функцией из math_converter."""
        return omml_to_mathml(omath_elem)

    def _table_to_html(self, table) -> str:
        """Преобразует таблицу из PPTX в HTML с сохранением форматирования (case, sub/sup, bold/italic) и формул."""
        rows_html = []
        for row_idx, row in enumerate(table.rows):
            cells_html = ""
            for cell in row.cells:
                # Собираем HTML ячейки из всех параграфов
                cell_parts = []
                for para in cell.text_frame.paragraphs:
                    cell_parts.append(self._paragraph_to_html(para))
                cell_content = " ".join(cell_parts)  # разделяем параграфы пробелом (как в монолите)
                tag = "th" if row_idx == 0 else "td"
                cells_html += f"<{tag}>{cell_content}</{tag}>"
            rows_html.append(f"<tr>{cells_html}</tr>")
        return f'<table class="data-table"><tbody>{"".join(rows_html)}</tbody></table>'

    def _paragraph_to_html(self, paragraph) -> str:
        """Преобразует параграф (из ячейки таблицы или фигуры) в HTML, сохраняя форматирование и формулы."""
        NS = {
            "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
            "a14": "http://schemas.microsoft.com/office/drawing/2014/main",
            "m": "http://schemas.openxmlformats.org/officeDocument/2006/math",
            "mc": "http://schemas.openxmlformats.org/markup-compatibility/2006",
        }
        p_element = paragraph._element
        parts = []
        for child in p_element:
            local = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if local == "r":
                omath = child.find(".//m:oMath", NS) or child.find(".//m:oMathPara", NS)
                if omath is None:
                    omath = child.find(".//mc:AlternateContent/mc:Choice//m:oMath", NS) or \
                            child.find(".//mc:AlternateContent/mc:Choice//m:oMathPara", NS)
                
                if omath is not None:
                    mathml = self._omml_to_mathml(omath)
                    if mathml:
                        parts.append(f'<span class="formula-container">{mathml}</span>')
                    else:
                        linear_parts = [t.text for t in omath.findall(".//m:t", NS) if t.text]
                        linear_text = "".join(linear_parts)
                        if linear_text:
                            parts.append(f'<span class="formula-fallback">{html.escape(linear_text)}</span>')
                    continue

                t_elems = child.findall(".//a:t", NS)
                txt = "".join(t.text or "" for t in t_elems)
                if not txt:
                    continue
                rPr = child.find(".//a:rPr", NS)
                sub = False
                sup = False
                bold = False
                italic = False
                if rPr is not None:
                    baseline = rPr.get("baseline")
                    if baseline:
                        try:
                            val = int(baseline)
                            if val > 0: sup = True
                            elif val < 0: sub = True
                        except: pass
                    
                    if not sub: sub = rPr.get("subscript") in ("1", "true", True)
                    if not sup: sup = rPr.get("superscript") in ("1", "true", True)
                    bold = rPr.get("b") in ("1", "true", True)
                    italic = rPr.get("i") in ("1", "true", True)
                
                content = html.escape(txt)
                if sub: content = f"<sub>{content}</sub>"
                elif sup: content = f"<sup>{content}</sup>"
                if bold: content = f"<strong>{content}</strong>"
                if italic: content = f"<em>{content}</em>"
                parts.append(content)

            elif local == "br":
                parts.append("<br/>")
            elif local in ("endParaRPr", "extLst"):
                continue
            else:
                omath = child.find(".//m:oMath", NS) or child.find(".//m:oMathPara", NS)
                if omath is not None:
                    mathml = self._omml_to_mathml(omath)
                    if mathml:
                        parts.append(f'<span class="formula-container">{mathml}</span>')
                    else:
                        linear_parts = [t.text for t in omath.findall(".//m:t", NS) if t.text]
                        linear_text = "".join(linear_parts)
                        if linear_text:
                            parts.append(f'<span class="formula-fallback">{html.escape(linear_text)}</span>')
        return "".join(parts)

    def _process_image(self, shape, slide_idx, slide_info):
        media_output_full = DESIGN_CONFIG['paths']['media_output_full']
        img_name = f"slide_{slide_idx + 1}_img_{len(slide_info['visuals']) + 1}.png"
        img_path_full = str(Path(media_output_full) / img_name)
        img_path_rel = f"{DESIGN_CONFIG['paths']['media_output']}/{img_name}"
        try:
            blob = shape.image.blob
            if _save_image_with_white_bg(blob, img_path_full):
                slide_info['visuals'].append({'src': img_path_rel, 'caption': '', 'pos': (shape.left, shape.top, shape.width, shape.height)})
                self.stats['images_ok'] += 1
        except Exception as e:
            logger.warning(f'Slide {slide_idx + 1}: Image fail: {e}')
            self.stats['images_fail'] += 1

    def _extract_captions_from_shapes(self, slide, slide_info: Dict[str, Any]) -> None:
        """Извлекает подписи к изображениям из текстовых блоков (порт из монолита)."""
        cfg = DESIGN_CONFIG.get("caption_search", {})
        vert_range_mm = cfg.get("vertical_range_mm", 15)
        max_gap_mm = cfg.get("max_gap_mm", 10)
        horiz_ratio = cfg.get("horizontal_overlap_ratio", 0.4)
        overlap_tol = cfg.get("overlap_tolerance", 0.5)
        priority = cfg.get("priority", "above")
        vert_range_emu = int(vert_range_mm * 36000)
        max_gap_emu = int(max_gap_mm * 36000)

        slide_num = slide_info.get("slide_num", "?")
        text_shapes_with_pos = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                txt = shape.text.strip()
                if txt and not _is_slide_number(txt) and len(txt) < 250:
                    text_shapes_with_pos.append({
                        "text": txt,
                        "left": shape.left,
                        "top": shape.top,
                        "width": shape.width,
                        "height": shape.height,
                    })

        used_texts = set()
        for vis in slide_info["visuals"]:
            v_left, v_top, v_w, v_h = vis["pos"]
            v_center_x = v_left + v_w / 2
            best_match = None
            best_score = 0.0
            second_best_score = 0.0

            # --- Pass 1: search ABOVE image ---
            search_top = v_top - vert_range_emu
            for ts in text_shapes_with_pos:
                if ts["text"] in used_texts:
                    continue
                ts_top = ts["top"]
                ts_bottom = ts["top"] + ts["height"]
                ts_center_x = ts["left"] + ts["width"] / 2
                ts_h = ts["height"]

                if ts_bottom < search_top:
                    continue

                v_gap = v_top - ts_bottom
                abs_gap = abs(v_gap)
                if abs_gap > max_gap_emu:
                    continue

                if v_gap < 0:
                    if ts_h > 0:
                        overlap_above = (v_top - ts_top) / ts_h
                        if overlap_above < overlap_tol:
                            continue
                    else:
                        continue

                h_offset = abs(ts_center_x - v_center_x)
                if h_offset > v_w * horiz_ratio:
                    continue
                
                h_thresh_score = v_w * 0.5 if v_w > 0 else 1
                h_score = max(0.0, 1.0 - (h_offset / h_thresh_score))
                v_score = max(0.0, 1.0 - (abs_gap / max_gap_emu)) if max_gap_emu > 0 else 0.0
                total_score = 0.7 * v_score + 0.3 * h_score

                if total_score > best_score:
                    second_best_score = best_score
                    best_score = total_score
                    best_match = ts
                elif total_score > second_best_score:
                    second_best_score = total_score

            # --- Pass 2: fallback BELOW image ---
            if priority == "above" and (best_match is None or best_score <= 0.6):
                v_bottom = v_top + v_h
                for ts in text_shapes_with_pos:
                    if ts["text"] in used_texts:
                        continue
                    ts_top = ts["top"]
                    ts_center_x = ts["left"] + ts["width"] / 2
                    h_offset = abs(ts_center_x - v_center_x)
                    if h_offset > v_w * horiz_ratio:
                        continue

                    v_gap_below = ts_top - v_bottom
                    if v_gap_below < 0 or v_gap_below > max_gap_emu:
                        continue

                    v_score = max(0.0, 1.0 - (v_gap_below / max_gap_emu)) if max_gap_emu > 0 else 0.0
                    h_thresh_score = v_w * 0.5 if v_w > 0 else 1
                    h_score = max(0.0, 1.0 - (h_offset / h_thresh_score))
                    total_score = 0.7 * v_score + 0.3 * h_score

                    if total_score > best_score:
                        second_best_score = best_score
                        best_score = total_score
                        best_match = ts
                    elif total_score > second_best_score:
                        second_best_score = total_score

            if best_match and best_score > 0.6:
                vis["caption"] = best_match["text"]
                used_texts.add(best_match["text"])
                
                def normalize(s: str) -> str:
                    s = re.sub(r"[^\w\s]", "", s.lower())
                    return " ".join(s.split()).strip()
                
                target = normalize(best_match["text"])
                if not target: continue

                new_items = []
                found_and_removed = False
                for item in slide_info.get("content_items", []):
                    if item["type"] == "text":
                        new_paras = []
                        for p in item["data"]:
                            p_plain = "".join(s for s in p if isinstance(s, str) and not s.startswith("[[["))
                            # Очищаем от внутренних маркеров форматирования [[...]] для правильного сопоставления
                            p_plain = re.sub(r'\[\[.*?\]\]', '', p_plain)
                            if normalize(p_plain) != target:
                                new_paras.append(p)
                            else:
                                found_and_removed = True
                        if new_paras:
                            item["data"] = new_paras
                            new_items.append(item)
                    else:
                        new_items.append(item)
                
                slide_info["content_items"] = new_items
                if best_score - second_best_score < 0.1 and second_best_score > 0:
                    logger.warning(f"Slide {slide_num}: ambiguous caption (score diff {best_score - second_best_score:.3f})")

    def _enrich_intro_slide(self, slide_info: Dict[str, Any]):
        """Извлекает имя докладчика и информацию о нем, имитируя логику монолита."""
        presenter_items = []
        for item in slide_info["content_items"]:
            if item["type"] == "text":
                for p in item["data"]:
                    # Очищаем от формул для поиска текста
                    txt = "".join(s for s in p if not s.startswith("[[[")).strip()
                    if txt:
                        presenter_items.append(txt)

        if presenter_items:
            title_candidate = presenter_items[0]
            speaker_candidate = ""
            info_candidate = ""
            remaining = presenter_items[1:]
            found_name = False

            for i, item in enumerate(remaining):
                pure_item = re.sub(r"\[\[.*?\]\]", "", item).strip()
                if pure_item.lower() in ("докладчик", "докладчик:", "выполнил", "автор"):
                    continue

                if not found_name:
                    name, info = self._clean_speaker_name("\n".join(remaining[i:]))
                    speaker_candidate = name
                    info_candidate = info
                    found_name = True
                    break

            if not found_name and remaining:
                name, info = self._clean_speaker_name("\n".join(remaining))
                speaker_candidate = name
                info_candidate = info

            slide_info["title"] = title_candidate
            slide_info["speaker_name"] = speaker_candidate if speaker_candidate else "Докладчик"
            slide_info["speaker_info"] = info_candidate
        else:
            slide_info["speaker_name"] = "Докладчик"
            slide_info["speaker_info"] = ""

    def _clean_speaker_name(self, text: str) -> tuple:
        """Вспомогательная функция для очистки имени докладчика (порт из монолита)."""
        lines = [l.strip() for l in text.split("\n") if l.strip()]
        if not lines:
            return "", ""
        
        name = lines[0]
        # Убираем "Докладчик:", "Автор:" и т.д.
        name = re.sub(r"^(Докладчик|Автор|Выполнил|Студент|Магистрант)[:\s]+", "", name, flags=re.IGNORECASE)
        
        info = "\n".join(lines[1:]) if len(lines) > 1 else ""
        return name, info

    def load_env(self) -> dict:
        """Загружает переменные окружения из .env файла."""
        env = {}
        path = DESIGN_CONFIG['paths']['env_path']
        if os.path.exists(path):
            with open(path, 'r', encoding='utf-8') as f:
                for line in f:
                    line = line.strip()
                    if '=' in line and (not line.startswith('#')):
                        k, v = line.split('=', 1)
                        env[k.strip()] = v.strip().strip('"').strip("'")
        return env

    def call_ai(self, prompt: str) -> str | None:
        """Вызывает LLM через OpenRouter API с кэшированием."""
        env = self.load_env()
        api_key = env.get('OPENROUTER_API_KEY')
        model = env.get('MODEL', 'openai/gpt-4o-mini')
        if not api_key:
            logger.warning('OPENROUTER_API_KEY не найден в .env')
            return None
        cache = LLMCache()
        cached = cache.get(prompt, model)
        if cached:
            logger.info(f'[LLM] Ответ получен из кэша для prompt: {prompt[:50]}...')
            return cached
        headers = {'Authorization': f'Bearer {api_key}', 'Content-Type': 'application/json'}
        payload = {'model': model, 'messages': [{'role': 'user', 'content': prompt}]}
        try:
            response = requests.post('https://openrouter.ai/api/v1/chat/completions', headers=headers, data=json.dumps(payload), timeout=30)
            if response.status_code == 200:
                result = response.json()
                content = result['choices'][0]['message']['content']
                cache.set(prompt, model, content)
                return content
            else:
                logger.error(f'Ошибка API: {response.status_code} - {response.text[:200]}')
        except Exception as e:
            logger.error(f'Исключение при вызове ИИ: {e}')
        return None

    def process_txt_files(self) -> None:
        """Обрабатывает вспомогательные текстовые файлы с использованием AI для структурирования."""
        files = {'Выводы.txt': {'type': 'conclusions', 'prompt': "Проанализируй текст и извлеки ключевые факты. СТРОГАЯ СТРУКТУРА: 1. Выполненные работы (МАКСИМУМ 6 пунктов): перечисли основные этапы и действия. 2. Результаты (МАКСИМУМ 6 пунктов): перечисли конкретные выводы, достижения и показатели. ПРАВИЛА: - СТРОГО соблюдай хронологический порядок событий. - Текст должен быть максимально сжатым, строгим и технически точным. - Объединяй весь опыт, не пропуская значимых деталей (даты, цифры), но формулируй их крайне емко. - Удали вводные слова, пояснения и 'воду'. ФОРМАТ: Верни только JSON-массив из кратких тезисов (сначала до 6 пунктов по работам, затем до 6 пунктов по результатам). Текст: {content}"}, 'Направление дальнейших исследований.txt': {'type': 'research', 'prompt': 'Извлеки основные направления дальнейших исследований. ПОРЯДОК: Сначала планируемые действия, затем ожидаемые эффекты и цели. ПРАВИЛА: - Соблюдай логическую и хронологическую последовательность. - Максимально краткий и сухой научный стиль. - Сохрани всю фактологическую базу данных. ФОРМАТ: Верни только JSON-массив строк. Текст: {content}'}}
        for filename, config in files.items():
            if not os.path.exists(filename):
                continue
            logger.info(f'Обработка текстового файла: {filename}')
            with open(filename, 'r', encoding='utf-8') as f:
                content = f.read()
            items = []
            ai_used = False
            api_key = self.load_env().get('OPENROUTER_API_KEY')
            if api_key:
                prompt = config['prompt'].format(content=content[:8000])
                ai_response = self.call_ai(prompt)
                if ai_response:
                    try:
                        cleaned = ai_response.strip()
                        if cleaned.startswith('```json'):
                            cleaned = cleaned.split('```json', 1)[1]
                        if cleaned.startswith('```'):
                            cleaned = cleaned.split('```', 1)[1]
                        if '```' in cleaned:
                            cleaned = cleaned.split('```', 1)[0]
                        cleaned = cleaned.strip()
                        parsed = json.loads(cleaned)
                        if isinstance(parsed, list) and all((isinstance(x, str) for x in parsed)):
                            items = [p.strip() for p in parsed if p.strip()]
                            ai_used = True
                            logger.info(f'[AI] Получено {len(items)} пунктов из {filename}')
                        else:
                            logger.warning(f'[AI] Ответ не является массивом строк, используем fallback')
                    except json.JSONDecodeError as e:
                        logger.warning(f'[AI] Не удалось распарсить JSON: {e}, используем fallback')
            if not items:
                paragraphs = [p.strip() for p in content.split('\n') if p.strip()]
                items = paragraphs
                if not ai_used:
                    logger.info(f'[Fallback] Использовано разбиение по строкам для {filename} ({len(items)} пунктов)')
            if config['type'] == 'conclusions':
                processed_items = []
                for item in items:
                    clean_item = clean_text(item)
                    icon = self.get_icon_by_text(clean_item)
                    processed_items.append(f"<div class='summary-item animate-up'><i data-lucide='{icon}' class='animate-marker' style='width: 1.6rem; height: 1.6rem; flex-shrink: 0;'></i> <div class='list-text'>{esc(clean_item)}</div></div>")
                mid = (len(processed_items) + 1) // 2
                left_html = ''.join(processed_items[:mid])
                right_html = ''.join(processed_items[mid:])
                self.slides_data.append({'title': 'Выводы и результаты этапа', 'layout_type': 'conclusions_dual', 'left_html': left_html, 'right_html': right_html, 'content_items': []})
            elif config['type'] == 'research':
                items_html = []
                for item in items:
                    clean_item = clean_text(item)
                    items_html.append(f"<div class='roadmap-item animate-up'><i data-lucide='rocket' class='animate-marker' style='width: 1.6rem; height: 1.6rem; flex-shrink: 0;'></i> <div class='list-text'>{esc(clean_item)}</div></div>")
                self.slides_data.append({'title': 'Направление дальнейших исследований', 'layout_type': 'research_roadmap', 'content_html': ''.join(items_html), 'content_items': []})

    def get_icon_by_text(self, text: str) -> str:
        """Возвращает идентификатор иконки на основе текста (по ключевым словам)."""
        text = text.lower()
        for icon, keywords in DESIGN_CONFIG['icon_mapping'].items():
            if any((kw in text for kw in keywords)):
                return icon
        return 'chevron-right'